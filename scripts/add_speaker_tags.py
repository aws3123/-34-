#!/usr/bin/env python3
"""为鎏金版PPT每页右上角添加汇报人标签"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PPT_INPUT  = r"D:\jineng\jinengshijain\jinengshijain\docs\ISO11820-答辩PPT-鎏金版.pptx"
PPT_OUTPUT = r"D:\jineng\jinengshijain\jinengshijain\docs\ISO11820-答辩PPT-鎏金版-团队汇报.pptx"

# ============================================================
# 汇报人分工（基于根目录计划书）
#
# 罗凌岚   — 骨架owner + 集成负责人（整体/架构/总结）
# chaoxingstar — 运行时owner（状态机/仿真/温控/数据采集）
# Nlpexi-e — UI owner（界面/交互/权限）
# userswdrs — 持久化owner（数据库/历史/校准/配置）
# Y-Yunye  — 导出owner（记录保存/CSV/导出/仿真框架）
# ============================================================
SPEAKER_MAP = {
    # 封面 + 大纲
    0:  "汇报人：罗凌岚",
    1:  "汇报人：罗凌岚",
    # Ch1 研究背景与意义
    2:  "汇报人：罗凌岚",     # 分隔页
    3:  "汇报人：罗凌岚",     # 研究背景
    4:  "汇报人：罗凌岚",     # 国内外现状
    5:  "汇报人：罗凌岚",     # 研究意义
    # Ch2 相关技术与理论
    6:  "汇报人：罗凌岚",     # 分隔页
    7:  "汇报人：chaoxingstar",  # 多阶段温控算法（运行时）
    8:  "汇报人：chaoxingstar",  # 稳定性判断（运行时）
    9:  "汇报人：罗凌岚",     # 技术栈（集成）
    10: "汇报人：罗凌岚",     # 通信与设计模式（架构）
    11: "汇报人：罗凌岚",     # 技术总结
    # Ch3 系统总体设计
    12: "汇报人：罗凌岚",     # 分隔页
    13: "汇报人：罗凌岚",     # 系统架构
    14: "汇报人：chaoxingstar",  # 状态机（运行时）
    15: "汇报人：userswdrs",     # 数据库设计（持久化）
    16: "汇报人：罗凌岚",     # 功能模块（总体）
    # Ch4 系统详细设计与实现
    17: "汇报人：罗凌岚",     # 分隔页
    18: "汇报人：chaoxingstar",  # 数据采集（DaqWorker）
    19: "汇报人：chaoxingstar",  # 试验流程控制（TestController）
    20: "汇报人：chaoxingstar",  # 温控算法实现（SensorSimulator）
    21: "汇报人：Nlpexi-e",      # 客户端界面（UI）
    22: "汇报人：Y-Yunye",       # 数据管理与导出（Export）
    23: "汇报人：Y-Yunye",       # 仿真模式（仿真框架）
    # Ch5 测试与总结
    24: "汇报人：罗凌岚",     # 测试与总结
    # 致谢
    25: "汇报人：罗凌岚",
}

# 深色背景页（标签用金色文字）
DARK_SLIDES = {0, 2, 6, 12, 17, 25}
# 标题栏为深色的内容页（标签放标题栏内右上角）
DARK_BAR_SLIDES = set(SPEAKER_MAP.keys()) - DARK_SLIDES


def add_speaker_tags():
    prs = Presentation(PPT_INPUT)

    for idx, speaker in SPEAKER_MAP.items():
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]

        # 确定标签颜色和位置
        if idx in DARK_SLIDES:
            # 深色背景页：金色文字，靠顶部
            text_color = RGBColor(0xC9, 0xA8, 0x4C)  # GOLD
            top = Inches(0.2)
        else:
            # 内容页：标签放在标题栏（深色）内右上角
            text_color = RGBColor(0xE8, 0xD5, 0xA3)  # GOLD_LIGHT
            top = Inches(0.45)

        # 添加小文字框
        tx = slide.shapes.add_textbox(
            Inches(10.5), top,
            Inches(2.5), Inches(0.35)
        )
        tf = tx.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = speaker
        p.font.size = Pt(10)
        p.font.color.rgb = text_color
        p.font.bold = False
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.RIGHT

    prs.save(PPT_OUTPUT)
    print(f"Done! Speaker tags added to {len(SPEAKER_MAP)} slides")
    print(f"Saved to: {PPT_OUTPUT}")

    # 验证
    prs2 = Presentation(PPT_OUTPUT)
    for idx, speaker in SPEAKER_MAP.items():
        if idx >= len(prs2.slides):
            continue
        slide = prs2.slides[idx]
        for s in slide.shapes:
            if hasattr(s, 'text') and '汇报人' in s.text:
                print(f"  Slide {idx+1:2d}: {s.text}")
                break


if __name__ == "__main__":
    add_speaker_tags()
