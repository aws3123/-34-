#!/usr/bin/env python3
"""生成数据库 ER 图并插入 PPT Slide 16 的占位位置"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

# ---- 颜色 ----
HEADER_BG   = (13, 43, 78)        # PRIMARY 深蓝
HEADER_FG   = (255, 255, 255)     # 白
BODY_BG     = (255, 255, 255)     # 白
BODY_BORDER = (189, 195, 208)     # SUBTLE_LINE
PK_COLOR    = (232, 141, 42)      # ACCENT 金色
FK_COLOR    = (39, 174, 96)       # ACCENT2 绿
FIELD_FG    = (52, 73, 94)        # TEXT_BODY
REL_LINE    = (26, 86, 158)       # SECONDARY
BG_COLOR    = (240, 243, 248)     # LIGHT_BG

# ---- 字体 ----
def get_font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except:
                pass
    return ImageFont.load_default()

FONT_HEADER = get_font(16, bold=True)
FONT_FIELD  = get_font(13)
FONT_TAG    = get_font(11, bold=True)
FONT_REL    = get_font(12)


# ---- 表定义 ----
tables = {
    "operators": {
        "fields": [
            ("id", "INTEGER PK", "pk"),
            ("username", "TEXT UNIQUE", ""),
            ("pwd", "TEXT (SHA256)", ""),
            ("role", "TEXT", ""),
            ("created_at", "TEXT", ""),
        ],
    },
    "apparatus": {
        "fields": [
            ("id", "INTEGER PK", "pk"),
            ("name", "TEXT", ""),
            ("model", "TEXT", ""),
            ("serial_number", "TEXT", ""),
            ("serial_port_config", "TEXT", ""),
            ("created_at", "TEXT", ""),
        ],
    },
    "productmaster": {
        "fields": [
            ("id", "INTEGER PK", "pk"),
            ("product_code", "TEXT UNIQUE", ""),
            ("test_id", "TEXT", ""),
            ("product_name", "TEXT", ""),
            ("specification", "TEXT", ""),
            ("height_mm", "REAL", ""),
            ("diameter_mm", "REAL", ""),
            ("created_at", "TEXT", ""),
        ],
    },
    "testmaster": {
        "fields": [
            ("productid", "TEXT PK/FK", "pk,fk"),
            ("testid", "TEXT PK", "pk"),
            ("testdate", "TEXT", ""),
            ("operator", "TEXT", ""),
            ("sample_name", "TEXT", ""),
            ("preweight", "REAL", ""),
            ("postweight", "REAL", ""),
            ("lostweight_per", "REAL", ""),
            ("deltatf", "REAL", ""),
            ("totaltesttime", "INTEGER", ""),
            ("flame_time", "INTEGER", ""),
            ("flame_duration", "INTEGER", ""),
            ("has_flame", "INTEGER", ""),
            ("env_temp", "REAL", ""),
            ("env_humidity", "REAL", ""),
            ("notes", "TEXT", ""),
            ("flag", "TEXT", ""),
        ],
    },
    "sensors": {
        "fields": [
            ("id", "INTEGER PK", "pk"),
            ("name", "TEXT", ""),
            ("type", "TEXT", ""),
            ("channel", "INTEGER", ""),
            ("range_scale", "TEXT", ""),
            ("created_at", "TEXT", ""),
        ],
    },
    "CalibrationRecords": {
        "fields": [
            ("id", "INTEGER PK", "pk"),
            ("sensor_id", "INTEGER FK", "fk"),
            ("calibration_date", "TEXT", ""),
            ("result_json", "TEXT", ""),
            ("technician", "TEXT", ""),
            ("notes", "TEXT", ""),
            ("operator", "TEXT", ""),
            ("created_at", "TEXT", ""),
        ],
    },
}

# ---- 关系 ----
relationships = [
    ("productmaster", "product_code", "testmaster", "productid", "1 : N"),
    ("sensors", "id", "CalibrationRecords", "sensor_id", "1 : N"),
    ("operators", "username", "testmaster", "operator", "1 : N"),
]


# ---- 绘图参数 ----
CARD_W = 260
ROW_H = 22
HEADER_H = 32
PADDING = 6
GAP_X = 80
GAP_Y = 60

def calc_table_h(fields):
    return HEADER_H + len(fields) * ROW_H + PADDING * 2


def draw_table(draw, x, y, name, fields):
    h = calc_table_h(fields)
    # 表体阴影
    draw.rounded_rectangle([x+3, y+3, x+CARD_W+3, y+h+3], radius=8, fill=(200, 200, 200))
    # 表体
    draw.rounded_rectangle([x, y, x+CARD_W, y+h], radius=8, fill=BODY_BG, outline=BODY_BORDER, width=2)
    # 表头
    draw.rounded_rectangle([x, y, x+CARD_W, y+HEADER_H], radius=8, fill=HEADER_BG)
    draw.rectangle([x, y+HEADER_H-8, x+CARD_W, y+HEADER_H], fill=HEADER_BG)
    # 表名
    bbox = draw.textbbox((0, 0), name, font=FONT_HEADER)
    tw = bbox[2] - bbox[0]
    draw.text((x + (CARD_W - tw) // 2, y + 7), name, fill=HEADER_FG, font=FONT_HEADER)
    # 字段
    fy = y + HEADER_H + PADDING
    for fname, ftype, tags in fields:
        tag_list = tags.split(",") if tags else []
        # PK/FK 标记
        prefix = ""
        color = FIELD_FG
        if "pk" in tag_list and "fk" in tag_list:
            prefix = "PK/FK "
            color = PK_COLOR
        elif "pk" in tag_list:
            prefix = "PK "
            color = PK_COLOR
        elif "fk" in tag_list:
            prefix = "FK "
            color = FK_COLOR

        line = f"{prefix}{fname}"
        draw.text((x + 12, fy), line, fill=color, font=FONT_FIELD)
        # 类型
        type_bbox = draw.textbbox((0, 0), ftype, font=FONT_TAG)
        type_w = type_bbox[2] - type_bbox[0]
        draw.text((x + CARD_W - type_w - 12, fy + 2), ftype, fill=(149, 165, 166), font=FONT_TAG)
        # 分隔线
        draw.line([(x + 10, fy + ROW_H - 1), (x + CARD_W - 10, fy + ROW_H - 1)],
                  fill=(236, 240, 241), width=1)
        fy += ROW_H
    return x, y, CARD_W, h


def find_field_center(table_pos, table_name, field_name):
    """找到某表某字段的中心 y 坐标和左右边界 x"""
    tx, ty, tw, _ = table_pos[table_name]
    fields = tables[table_name]["fields"]
    fy = ty + HEADER_H + PADDING
    for fname, _, _ in fields:
        if fname == field_name:
            return tx, fy + ROW_H // 2, tx + tw
        fy += ROW_H
    return tx, ty + HEADER_H // 2, tx + tw


def draw_relationship(draw, table_pos, t1, f1, t2, f2, label):
    x1, y1, r1 = find_field_center(table_pos, t1, f1)
    x2, y2, r2 = find_field_center(table_pos, t2, f2)
    # 选择最近的连接边
    if r1 < x2:
        sx, ex = r1, x2
    elif x1 > r2:
        sx, ex = x1, r2
    else:
        sx, ex = r1, x2

    # 绘制折线
    mid_x = (sx + ex) // 2
    draw.line([(sx, y1), (mid_x, y1)], fill=REL_LINE, width=2)
    draw.line([(mid_x, y1), (mid_x, y2)], fill=REL_LINE, width=2)
    draw.line([(mid_x, y2), (ex, y2)], fill=REL_LINE, width=2)

    # 关系标签
    bbox = draw.textbbox((0, 0), label, font=FONT_REL)
    lw = bbox[2] - bbox[0]
    lx = mid_x - lw // 2
    ly = (y1 + y2) // 2 - 10
    draw.rounded_rectangle([lx-4, ly-2, lx+lw+4, ly+16], radius=4, fill=BG_COLOR)
    draw.text((lx, ly), label, fill=REL_LINE, font=FONT_REL)


def generate_er_diagram(output_path):
    # 布局: 3行2列
    # Row 0: operators, apparatus
    # Row 1: productmaster, testmaster (center, large)
    # Row 2: sensors, CalibrationRecords
    margin = 40
    col0_x = margin
    col1_x = margin + CARD_W + GAP_X

    positions = {}
    row0_y = margin
    positions["operators"] = (col0_x, row0_y)
    positions["apparatus"] = (col1_x, row0_y)

    row1_y = row0_y + calc_table_h(tables["operators"]["fields"]) + GAP_Y
    # testmaster 很大，放右边
    positions["productmaster"] = (col0_x, row1_y)
    positions["testmaster"] = (col1_x, row1_y)

    row2_y = row1_y + calc_table_h(tables["testmaster"]["fields"]) + GAP_Y
    positions["sensors"] = (col0_x, row2_y)
    positions["CalibrationRecords"] = (col1_x, row2_y)

    # 计算画布大小
    max_x = max(x + CARD_W for x, y in positions.values()) + margin
    max_y = max(y + calc_table_h(tables[name]["fields"]) for name, (x, y) in positions.items()) + margin

    img = Image.new("RGB", (max_x, max_y), BG_COLOR)
    draw = ImageDraw.Draw(img)

    # 绘制所有表
    table_pos = {}
    for name, (x, y) in positions.items():
        info = draw_table(draw, x, y, name, tables[name]["fields"])
        table_pos[name] = info

    # 绘制关系线
    for t1, f1, t2, f2, label in relationships:
        draw_relationship(draw, table_pos, t1, f1, t2, f2, label)

    # 添加标题
    title = "ISO 11820 System Database ER Diagram"
    bbox = draw.textbbox((0, 0), title, font=FONT_HEADER)
    tw = bbox[2] - bbox[0]
    # 不需要标题，PPT里已有

    img.save(output_path, "PNG", dpi=(150, 150))
    print(f"ER diagram saved to: {output_path}")
    print(f"Image size: {img.size}")
    return output_path


def insert_into_ppt(ppt_path, img_path):
    """将 ER 图插入 PPT Slide 16 的占位位置"""
    prs = Presentation(ppt_path)
    slide = prs.slides[15]  # 第16页 (0-indexed)

    # 占位位置: Inches(9.5), Inches(2.15), Inches(3.2), Inches(3.0)
    left = Inches(9.3)
    top = Inches(2.15)
    width = Inches(3.6)
    height = Inches(3.0)

    # 删除占位框（最后两个含"此处需要"文字的shape）
    shapes_to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and '此处需要插上数据库ER图' in shape.text:
            shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)

    # 插入图片
    slide.shapes.add_picture(img_path, left, top, width, height)

    prs.save(ppt_path)
    print(f"ER diagram inserted into PPT Slide 16: {ppt_path}")


if __name__ == "__main__":
    docs_dir = r"D:\jineng\jinengshijain\jinengshijain\docs"
    img_path = os.path.join(docs_dir, "er_diagram.png")
    ppt_path = os.path.join(docs_dir, "ISO11820-答辩PPT-项目版.pptx")
    ppt_output = os.path.join(docs_dir, "ISO11820-答辩PPT-项目版-含ER图.pptx")

    generate_er_diagram(img_path)
    insert_into_ppt(ppt_path, img_path)
    print("Done!")
