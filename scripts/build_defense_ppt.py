#!/usr/bin/env python3
"""ISO 11820 建筑材料不燃性实验自动化测控系统 — 答辩PPT生成器
模仿模板PPT的章节结构（封面+目录+5章+致谢 = 26页）"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# 配色方案 - 学术蓝
# ============================================================
PRIMARY = RGBColor(0x0D, 0x2B, 0x4E)
SECONDARY = RGBColor(0x1A, 0x56, 0x9E)
ACCENT = RGBColor(0xE8, 0x8D, 0x2A)
ACCENT2 = RGBColor(0x27, 0xAE, 0x60)
LIGHT_BG = RGBColor(0xF0, 0xF3, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
TEXT_BODY = RGBColor(0x34, 0x49, 0x5E)
SUBTLE_LINE = RGBColor(0xBD, 0xC3, 0xD0)
CARD_BG = RGBColor(0xFA, 0xFB, 0xFD)
RED_ALERT = RGBColor(0xE7, 0x4C, 0x3C)

BODY_FONT = 'Microsoft YaHei'
EN_FONT = 'Arial'

# 标准16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class PPTBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.slide_num = 0

    # ---- 基础绘图 ----
    def _set_bg(self, slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _add_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_rounded_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_textbox(self, slide, left, top, width, height, text, font_size=Pt(18),
                     color=TEXT_BODY, bold=False, align=PP_ALIGN.LEFT, font_name=BODY_FONT,
                     line_spacing=1.3, auto_fit=False):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        if auto_fit:
            txBox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        return txBox

    def _add_multiline_textbox(self, slide, left, top, width, height, lines,
                               default_size=Pt(18), default_color=TEXT_BODY,
                               default_bold=False, line_spacing=1.3, auto_fit=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        if auto_fit:
            txBox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf = txBox.text_frame
        for i, line in enumerate(lines):
            if isinstance(line, str):
                text, size, color, bold = line, default_size, default_color, default_bold
            else:
                text = line[0]
                size = line[1] if len(line) > 1 else default_size
                color = line[2] if len(line) > 2 else default_color
                bold = line[3] if len(line) > 3 else default_bold
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = BODY_FONT
            p.line_spacing = line_spacing
            p.space_after = Pt(3)
            p.space_before = Pt(0)
        return txBox

    def _page_number(self, slide):
        self.slide_num += 1
        self._add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35),
                          str(self.slide_num), Pt(11), SUBTLE_LINE, align=PP_ALIGN.RIGHT)

    def _title_bar(self, slide, title, subtitle=None):
        self._add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.35), PRIMARY)
        self._add_rect(slide, Inches(0), Inches(1.35), SLIDE_W, Inches(0.04), ACCENT)
        self._add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.85),
                          title, Pt(30), WHITE, bold=True)
        if subtitle:
            self._add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.4),
                              subtitle, Pt(14), RGBColor(0xBB, 0xCC, 0xDD))

    def _content_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, LIGHT_BG)
        self._title_bar(slide, title, subtitle)
        self._page_number(slide)
        return slide

    def _chapter_divider(self, number, title_cn, title_en):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, LIGHT_BG)
        # 左侧装饰条
        self._add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, PRIMARY)
        # 大数字
        self._add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(2),
                          number, Pt(96), PRIMARY, bold=True, font_name=EN_FONT)
        # 横线
        self._add_rect(slide, Inches(0.8), Inches(3.8), Inches(3), Inches(0.04), ACCENT)
        # 中文标题
        self._add_textbox(slide, Inches(0.8), Inches(4.1), Inches(8), Inches(0.8),
                          title_cn, Pt(36), PRIMARY, bold=True)
        # 英文标题
        self._add_textbox(slide, Inches(0.8), Inches(4.9), Inches(8), Inches(0.5),
                          title_en, Pt(16), SUBTLE_LINE, font_name=EN_FONT)
        # 右侧竖线装饰
        self._add_rect(slide, Inches(12.5), Inches(1.0), Inches(0.03), Inches(5.5), SECONDARY)
        self._page_number(slide)
        return slide

    def _card(self, slide, left, top, width, height, icon_text, title, body_lines, color=SECONDARY):
        self._add_rounded_rect(slide, left, top, width, height, WHITE)
        self._add_rect(slide, left, top, width, Inches(0.06), color)
        self._add_textbox(slide, left + Inches(0.25), top + Inches(0.15), Inches(0.5), Inches(0.45),
                          icon_text, Pt(22), color, bold=True)
        self._add_textbox(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.5),
                          Inches(0.35), title, Pt(14), TEXT_DARK, bold=True)
        self._add_multiline_textbox(slide, left + Inches(0.25), top + Inches(0.95),
                                    width - Inches(0.5), height - Inches(1.1),
                                    body_lines, default_size=Pt(11), default_color=TEXT_BODY,
                                    line_spacing=1.35, auto_fit=True)

    def _bullet_list(self, slide, left, top, width, height, items, size=Pt(15),
                     color=TEXT_BODY, spacing=1.4):
        lines = [f"  {item}" if isinstance(item, str) else f"  {item[0]}" for item in items]
        self._add_multiline_textbox(slide, left, top, width, height, lines,
                                    default_size=size, default_color=color,
                                    line_spacing=spacing, auto_fit=True)

    def _numbered_item(self, slide, left, top, width, num, title, desc, color=SECONDARY):
        badge = self._add_rounded_rect(slide, left, top, Inches(0.42), Inches(0.42), color)
        self._add_textbox(slide, left, top + Inches(0.02), Inches(0.42), Inches(0.38),
                          str(num), Pt(18), WHITE, bold=True, align=PP_ALIGN.CENTER)
        self._add_textbox(slide, left + Inches(0.55), top, width - Inches(0.55), Inches(0.3),
                          title, Pt(15), TEXT_DARK, bold=True)
        self._add_textbox(slide, left + Inches(0.55), top + Inches(0.32), width - Inches(0.55),
                          Inches(0.5), desc, Pt(11), TEXT_BODY, auto_fit=True)

    def _add_table(self, slide, left, top, col_widths, headers, rows, header_color=PRIMARY, row_h=Inches(0.5)):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        total_w = sum(col_widths)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, row_h * n_rows)
        table = table_shape.table
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw
        for ci, h in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = str(val)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF1)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = TEXT_BODY
                    p.font.name = BODY_FONT
                    p.alignment = PP_ALIGN.CENTER
        return table_shape

    def _flow_arrow(self, slide, left, top, width, text, color=SECONDARY):
        shape = self._add_rounded_rect(slide, left, top, width, Inches(0.65), color)
        self._add_textbox(slide, left, top + Inches(0.04), width, Inches(0.57),
                          text, Pt(10), WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.1)

    def _arrow_right(self, slide, left, top):
        self._add_textbox(slide, left, top + Inches(0.08), Inches(0.25), Inches(0.45),
                          "→", Pt(18), ACCENT, bold=True, align=PP_ALIGN.CENTER)

    def _image_placeholder(self, slide, left, top, width, height, hint):
        """添加图片占位指示框 — 提醒在此处插入截图"""
        PH_BORDER = RGBColor(0xE7, 0x4C, 0x3C)
        PH_BG = RGBColor(0xFF, 0xF3, 0xCD)
        PH_TEXT = RGBColor(0xC0, 0x39, 0x2B)
        shape = self._add_rounded_rect(slide, left, top, width, height, PH_BG)
        shape.line.color.rgb = PH_BORDER
        shape.line.width = Pt(2.5)
        # 相机图标 + 提示文字
        self._add_textbox(slide, left, top + Inches(0.08), width, Inches(0.35),
                          "📷", Pt(16), PH_BORDER, bold=True, align=PP_ALIGN.CENTER)
        self._add_textbox(slide, left + Inches(0.1), top + Inches(0.4),
                          width - Inches(0.2), height - Inches(0.5),
                          hint, Pt(11), PH_TEXT, bold=True,
                          align=PP_ALIGN.CENTER, line_spacing=1.4)


# ============================================================
def build():
    b = PPTBuilder()

    # ================================================================
    # SLIDE 1: 封面 (对应模板Slide 1)
    # ================================================================
    slide = b.prs.slides.add_slide(b.prs.slide_layouts[6])
    b._set_bg(slide, PRIMARY)
    b._add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    b._add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x08, 0x1F, 0x3A))
    b._add_textbox(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.5),
                   "建筑材料不燃性实验自动化测控系统", Pt(44), WHITE, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(3.3), Inches(11), Inches(0.6),
                   "研究与实现", Pt(36), ACCENT, bold=False)
    b._add_rect(slide, Inches(1.2), Inches(4.2), Inches(3), Inches(0.04), ACCENT)
    b._add_textbox(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(0.5),
                   "—— 毕业论文答辩", Pt(22), RGBColor(0xBB, 0xCC, 0xDD))
    b._add_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(0.5),
                   ".NET 8  ·  WinForms  ·  SQLite  ·  OxyPlot  ·  EPPlus  ·  MathNet.Numerics",
                   Pt(14), RGBColor(0x99, 0xAA, 0xBB))
    b._add_textbox(slide, Inches(1.2), Inches(6.0), Inches(5), Inches(0.4),
                   "答辩人：罗凌岚  |  指导教师：肖斌教授  |  2026年5月", Pt(14), RGBColor(0x88, 0x99, 0xAA))
    b.slide_num += 1

    # ================================================================
    # SLIDE 2: 目录 (对应模板Slide 2)
    # ================================================================
    slide = b._content_slide("汇报大纲", "CONTENTS")
    agenda_items = [
        ("01", "研究背景与意义", "研究背景、国内外现状、研究目标"),
        ("02", "相关技术与理论", "温控算法、稳定性判断、开发技术"),
        ("03", "系统总体设计", "架构设计、状态机、数据库、功能模块"),
        ("04", "系统详细设计与实现", "服务端实现、客户端实现、仿真模式"),
        ("05", "系统测试与总结", "功能测试、实验数据、总结与展望"),
    ]
    for i, (num, title, desc) in enumerate(agenda_items):
        col = i % 3
        row = i // 3
        left = Inches(0.6 + col * 4.2)
        top = Inches(2.0 + row * 2.2)
        b._numbered_item(slide, left, top, Inches(3.8), num, title, desc)

    # ================================================================
    # SLIDE 3: Chapter 1 分隔页 — 研究背景与意义
    # ================================================================
    b._chapter_divider("01", "研究背景与意义", "BACKGROUND AND SIGNIFICANCE")

    # ================================================================
    # SLIDE 4: 1.1 研究背景 (对应模板Slide 4)
    # ================================================================
    slide = b._content_slide("1. 研究背景与意义 · 研究背景", "Chapter 1 · 研究背景与意义")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "研究背景", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "建筑材料防火安全性能是构筑社会公共安全防线的基石",
        "不燃性是评价建筑材料防火等级最基础、最重要的性能指标",
        "ISO 11820-2020 标准规定：750°C±5°C 高温条件下连续稳定10分钟",
        "传统人工测试模式效率低、过程可控性不强、数据一致性不佳",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "传统人工测试的痛点", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "操作人员需全程密切关注并手动调节，劳动强度极大",
        "凭经验判断温度稳定，难以实现 ±0.5°C 的高精度控制",
        "人工记录和计算数据效率低，容易引入主观误差",
        "无法满足 ISO 11820-2020 对过程控制精度的严格条件",
    ], Pt(13))

    b._add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.0), WHITE)
    b._add_rect(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.06), ACCENT)
    b._add_textbox(slide, Inches(1.1), Inches(4.95), Inches(2), Inches(0.35),
                   "研究目标", Pt(18), ACCENT, bold=True)
    b._add_textbox(slide, Inches(1.1), Inches(5.4), Inches(10.8), Inches(1.2),
                   "开发一套严格依照 ISO 11820-2020 国际标准、能替代人工操作达成全流程自动化控制的测控系统，"
                   "以大幅提升建筑材料不燃性测试的效率、准确性以及数据规范性。",
                   Pt(14), TEXT_BODY, line_spacing=1.5)

    # ================================================================
    # SLIDE 5: 1.2 国内外研究现状 (对应模板Slide 5)
    # ================================================================
    slide = b._content_slide("1. 研究背景与意义 · 国内外现状", "Chapter 1 · 研究背景与意义")
    b._card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0),
            "🇨🇳", "国内研究现状", [
                "GB/T 5464-2010 标准为实验搭建了基础框架",
                "自动化程度较低，主要依赖人工操作",
                "有学者探索测试设备改进：自动定位传感器、自动测量",
                "模块化设计提升了试验稳定性",
                "但在 C# 等高级语言开发的智能测控系统方面应用案例较少",
                "难以满足实时数据分析和远程监控的需求",
            ])
    b._card(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(5.0),
            "🌍", "国外研究现状", [
                "欧美主导，聚焦智能化与自动化技术深度整合",
                "数据驱动方法突显自动化系统在提升测试精度中的必要性",
                "AI/ML/DL 技术在建筑4.0中广泛应用",
                "NIST 提出融合红外热成像与气体传感器的动态监测平台",
                "但缺乏针对 ISO 11820 标准流程的定制化开发案例",
                "面临高成本和数据兼容性挑战",
            ])

    # ================================================================
    # SLIDE 6: 1.3 研究意义 (对应模板Slide 6)
    # ================================================================
    slide = b._content_slide("1. 研究背景与意义 · 研究意义", "Chapter 1 · 研究背景与意义")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
                   "本研究的核心价值", Pt(20), PRIMARY, bold=True)

    values = [
        ("🔬", "技术层面", "为类似标准检测方法的自动化改造\n提供可复用的技术框架和工程范例", SECONDARY),
        ("📊", "数据层面", "标准化、结构化数据为实验室信息化\n管理与材料性能大数据分析奠定基础", ACCENT2),
        ("⚡", "效率层面", "试验效率提升3-5倍，人工干预时间\n从3-4小时减少到约10分钟", ACCENT),
        ("🏗", "行业层面", '推动检测行业从"经验判断"朝\n"数据驱动"转型升级', PRIMARY),
    ]
    for i, (icon, title, desc, color) in enumerate(values):
        left = Inches(0.5 + i * 3.15)
        b._card(slide, left, Inches(2.3), Inches(2.9), Inches(2.5), icon, title, [desc], color)

    b._add_rounded_rect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), PRIMARY)
    b._add_textbox(slide, Inches(1.2), Inches(5.35), Inches(10.5), Inches(0.35),
                   "研究方法", Pt(18), ACCENT, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(5.8), Inches(10.5), Inches(0.7),
                   "运用工业测控技术和现代软件工程方法，基于 .NET 8 框架与 WinForms 三层架构，"
                   "集成状态机模式、多阶段自适应温控算法和线性拟合稳定性判断算法，"
                   "设计并实现一套针对建筑材料不燃性试验的自动化测控系统。",
                   Pt(13), WHITE, line_spacing=1.5)

    # ================================================================
    # SLIDE 7: Chapter 2 分隔页 — 相关技术与理论
    # ================================================================
    b._chapter_divider("02", "相关技术与理论", "KEY TECHNOLOGIES AND THEORIES")

    # ================================================================
    # SLIDE 8: 2.1 多阶段自适应温控算法 (对应模板Slide 8)
    # ================================================================
    slide = b._content_slide("2. 关键技术 · 多阶段自适应温控算法", "Chapter 2 · 相关技术与理论")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "算法核心思想", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        '采用"分段逼近，平滑过渡"的控制策略',
        "升温过程划分为五个逻辑阶段",
        "各阶段运用不同的功率输出策略",
        "融合开环控制响应快与闭环控制精度高的优势",
    ], Pt(13))

    # 功率阶梯表
    stages = [
        ["温度区间", "功率策略", "控制方式"],
        ["室温 ~ 300°C", "30% 额定功率（软启动）", "开环"],
        ["300°C ~ 500°C", "50% 额定功率", "开环"],
        ["500°C ~ 600°C", "70% 额定功率", "开环"],
        ["600°C ~ 700°C", "90% 额定功率", "开环"],
        ["> 700°C", "切换 PID 闭环", "闭环"],
    ]
    b._add_table(slide, Inches(7.0), Inches(1.7),
                 [Inches(2.0), Inches(2.2), Inches(1.5)], stages[0], stages[1:], row_h=Inches(0.45))

    b._card(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(1.5),
            "🔑", "模式切换", [
                "当炉温突破 700°C 阈值后执行关键模式切换",
                "通过 Modbus 协议向 PID 控制器写入目标温度 750.0°C",
                "将控制权移交至硬件 PID 闭环控制器",
            ])
    b._card(slide, Inches(7.0), Inches(4.5), Inches(5.6), Inches(1.5),
            "✅", "算法优势", [
                "最大超调量仅 2.1°C（全功率方案为 32.6°C）",
                "稳态波动范围 ±0.5°C 以内",
                "升温速度与超调抑制之间取得最优平衡",
            ], ACCENT2)
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.75),
                         "（此处需要插上多阶段升温过程功率阶梯示意图 —— 对应论文图2-1）")

    # ================================================================
    # SLIDE 9: 2.2 线性拟合稳定性判断算法 (对应模板Slide 9)
    # ================================================================
    slide = b._content_slide("2. 关键技术 · 温度稳定性自动判断", "Chapter 2 · 相关技术与理论")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "数学基础：最小二乘线性回归", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "动态滑动窗口：存储最近10分钟温度数据（600个采样点）",
        "采样周期 800ms，每分钟调用 CheckStartCriteria()",
        "使用 MathNet.Numerics 库的 Fit.Line() 函数",
        "计算拟合直线斜率 β，得到温度漂移量 ΔDrift = |β| × 600",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "三重条件联合判定", Pt(20), PRIMARY, bold=True)
    conditions = [
        ["判据", "条件", "ISO标准要求"],
        ["温度范围", "T ∈ [745°C, 755°C]", "目标温度 ±5°C"],
        ["温度漂移", "ΔDrift ≤ 2.0°C/10min", "≤ 2°C/10min"],
        ["温度偏差", "Tmax - Tavg ≤ 1.0°C", "≤ 1°C"],
    ]
    b._add_table(slide, Inches(7.0), Inches(2.15),
                 [Inches(1.5), Inches(2.5), Inches(1.8)], conditions[0], conditions[1:], row_h=Inches(0.45))

    b._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "🔄", "持续性验证机制", [
                "每项条件维护独立倒计时计数器（初始600）",
                "条件满足则计数器减1，否则重置为600",
                "三个计数器同时归零才判定稳定",
                "杜绝因瞬时扰动导致的误判",
            ])
    b._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "📈", "线性拟合 vs 其他方法", [
                "极差法：对单次野值极其敏感",
                "方差法：无法识别缓慢单调漂移",
                "线性拟合：抗干扰性 + 趋势敏感性",
                "最小二乘对全部数据点做平均化处理",
            ])
    b._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(1.5),
            "🎯", "核心优势", [
                "将ISO标准定性要求转化为精确数学模型",
                "全自动、客观、连续验证",
            ], ACCENT2)
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(0.8), Inches(6.15), Inches(5.6), Inches(0.75),
                         "（此处需要插上线性拟合算法核心代码截图 —— 对应论文图5-2）")
    b._image_placeholder(slide, Inches(6.8), Inches(6.15), Inches(5.9), Inches(0.75),
                         "（此处需要插上三重条件判断核心代码截图 —— 对应论文图5-3）")

    # ================================================================
    # SLIDE 10: 2.3 开发技术栈 (对应模板Slide 10)
    # ================================================================
    slide = b._content_slide("2. 关键技术 · 开发技术栈", "Chapter 2 · 相关技术与理论")
    techs = [
        ("🖥", ".NET 8 / C# 10", ["跨平台高性能框架", "WinForms 桌面客户端", "异步编程模型"]),
        ("📊", "OxyPlot 2.1.2", ["实时温度曲线绘制", "10分钟滑动窗口", "自适应数据抽稀"]),
        ("🗄", "SQLite + EF Core", ["轻量级嵌入式数据库", "6张核心业务表", "ORM 对象关系映射"]),
        ("📄", "EPPlus + PDFsharp", ["Excel 格式化报表", "PDF 正式报告", "CSV 原始数据"]),
        ("🧮", "MathNet.Numerics", ["线性回归分析", "温度漂移计算", "数值计算库"]),
        ("📝", "Serilog 4.0", ["结构化日志记录", "文件滚动存储", "故障追踪审计"]),
    ]
    for i, (icon, title, lines) in enumerate(techs):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.7 + row * 2.8)
        b._card(slide, left, top, Inches(3.9), Inches(2.5), icon, title, lines)

    # ================================================================
    # SLIDE 11: 2.4 Modbus RTU 与状态机 (对应模板Slide 11)
    # ================================================================
    slide = b._content_slide("2. 关键技术 · 通信协议与设计模式", "Chapter 2 · 相关技术与理论")
    b._card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(2.5),
            "📡", "Modbus RTU 通信协议", [
                "工业标准串行通信协议，广泛应用于温控设备",
                "功能码 0x03：读取寄存器（温度数据 0x0102）",
                "功能码 0x06：写入寄存器（目标温度、控制模式）",
                "FluentModbus 库封装，超时重试 + 异常处理",
            ])
    b._card(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(2.5),
            "⚙", "状态机设计模式", [
                "将复杂试验流程离散为有限个状态",
                "状态间转换由明确定义的事件触发",
                "确保试验严格按 ISO 11820 标准步骤执行",
                "代码结构清晰，可维护性好，易于扩展",
            ])

    b._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "🏗", "三层架构设计", [
                "表示层：WinForms UI 界面",
                "业务逻辑层：状态机 + 算法",
                "数据访问层：SQLite + EF Core",
                "高内聚、低耦合",
            ])
    b._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "🔒", "线程安全机制", [
                "后台线程驱动数据采集（800ms）",
                "MainForm 通过 InvokeRequired 检测",
                "Invoke 安全调度 UI 更新",
                "RuntimeSnapshot 不可变快照",
            ])
    b._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(2.3),
            "🔌", "仿真测试支持", [
                "虚拟串口 + Modbus Slave 模拟器",
                "无需物理硬件即可完整测试",
                "分层可配置仿真开关",
                "极大降低开发验证成本",
            ], ACCENT2)

    # ================================================================
    # SLIDE 12: 2.5 技术总结 (对应模板Slide 12)
    # ================================================================
    slide = b._content_slide("2. 关键技术 · 核心技术体系", "Chapter 2 · 相关技术与理论")
    # 技术体系总览
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
                   "核心技术体系总览", Pt(20), PRIMARY, bold=True)

    core_techs = [
        ("多阶段自适应\n温控算法", SECONDARY, "室温→750°C\n五阶段功率控制\n+ PID 闭环切换"),
        ("线性拟合\n稳定性判断", ACCENT2, "MathNet.Numerics\n三重条件联合判定\n600秒持续验证"),
        ("状态机\n流程控制", PRIMARY, "5状态精确建模\n自动迁移 + 回退\nISO标准严格执行"),
        ("Modbus RTU\n硬件通信", ACCENT, "PID温控器通信\n传感器数据采集\n超时重试容错"),
        ("数据持久化\n与报表", RGBColor(0x8E, 0x44, 0xAD), "SQLite 6张表\nCSV/Excel/PDF\n自动化报告生成"),
    ]
    for i, (title, color, desc) in enumerate(core_techs):
        left = Inches(0.3 + i * 2.6)
        b._add_rounded_rect(slide, left, Inches(2.3), Inches(2.4), Inches(1.6), color)
        b._add_textbox(slide, left + Inches(0.15), Inches(2.4), Inches(2.1), Inches(0.65),
                       title, Pt(13), WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.2)
        b._add_textbox(slide, left + Inches(0.15), Inches(3.1), Inches(2.1), Inches(0.7),
                       desc, Pt(10), RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER, line_spacing=1.3)

    b._add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.35),
                   "技术选型理由", Pt(18), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(4.75), Inches(5.8), Inches(0.5), [
        ".NET 8：统一平台，高性能运行时，跨平台支持",
        "WinForms：成熟的桌面 GUI 框架，控件丰富，硬件交互稳定",
        "SQLite：轻量嵌入式，无需独立服务，适合单机部署",
        "OxyPlot：高性能图表库，支持实时数据流绘制",
    ], Pt(12))
    b._bullet_list(slide, Inches(7.0), Inches(4.75), Inches(5.8), Inches(0.5), [
        "MathNet.Numerics：权威数值计算库，线性回归精确可靠",
        "EPPlus/PDFsharp：成熟的 Office 文档生成方案",
        "Serilog：结构化日志，支持多目标输出",
        "FluentModbus：简洁的 Modbus 协议封装",
    ], Pt(12))

    # ================================================================
    # SLIDE 13: Chapter 3 分隔页 — 系统总体设计
    # ================================================================
    b._chapter_divider("03", "系统总体设计", "SYSTEM ARCHITECTURE DESIGN")

    # ================================================================
    # SLIDE 14: 3.1 系统架构设计 (对应模板Slide 14)
    # ================================================================
    slide = b._content_slide("3. 系统总体设计 · 系统架构", "Chapter 3 · 系统总体设计")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.35),
                   "四层技术架构", Pt(20), PRIMARY, bold=True)

    layers = [
        ("表示层 (UI)", "WinForms 图形界面 · 温度曲线 · 状态显示 · 操作控制", RGBColor(0x34, 0x98, 0xDB)),
        ("业务逻辑层", "TestController 状态机 · 温控算法 · 稳定性判断 · 广播机制", SECONDARY),
        ("数据访问层", "SQLite + EF Core · 数据持久化 · CSV/Excel/PDF 导出", RGBColor(0x8E, 0x44, 0xAD)),
        ("硬件抽象层", "Modbus RTU · PID控制器 · 传感器采集 · 仿真引擎", PRIMARY),
    ]
    for i, (name, desc, color) in enumerate(layers):
        top = Inches(2.2 + i * 1.0)
        b._add_rounded_rect(slide, Inches(0.8), top, Inches(11.5), Inches(0.82), color)
        b._add_textbox(slide, Inches(1.1), top + Inches(0.08), Inches(3), Inches(0.35),
                       name, Pt(16), WHITE, bold=True)
        b._add_textbox(slide, Inches(4.3), top + Inches(0.1), Inches(7.5), Inches(0.35),
                       desc, Pt(12), RGBColor(0xEE, 0xEE, 0xEE))

    b._add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.35),
                   "设计原则：Core 零外部依赖 | UI 不写业务逻辑 | 共享契约跨层交互 | 模块高内聚低耦合",
                   Pt(12), ACCENT, bold=True, align=PP_ALIGN.CENTER)
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(7.5), Inches(5.0), Inches(5.0), Inches(1.6),
                         "（此处需要插上系统四层架构图 —— 对应论文图4-1）")

    # ================================================================
    # SLIDE 15: 3.2 状态机设计 (对应模板Slide 15)
    # ================================================================
    slide = b._content_slide("3. 系统总体设计 · 状态机设计", "Chapter 3 · 系统总体设计")
    states = [
        ("Idle", "空闲", RGBColor(0x95, 0xA5, 0xA6)),
        ("Preparing", "升温中", SECONDARY),
        ("Ready", "就绪", ACCENT2),
        ("Recording", "记录中", ACCENT),
        ("Complete", "完成", PRIMARY),
    ]
    for i, (code, label, color) in enumerate(states):
        left = Inches(0.5 + i * 2.55)
        b._add_rounded_rect(slide, left, Inches(1.9), Inches(2.2), Inches(0.95), color)
        b._add_textbox(slide, left, Inches(1.98), Inches(2.2), Inches(0.42),
                       code, Pt(18), WHITE, bold=True, align=PP_ALIGN.CENTER)
        b._add_textbox(slide, left, Inches(2.4), Inches(2.2), Inches(0.35),
                       label, Pt(12), RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)
        if i < len(states) - 1:
            b._add_textbox(slide, left + Inches(2.25), Inches(2.1), Inches(0.3), Inches(0.4),
                           "→", Pt(26), ACCENT, bold=True, align=PP_ALIGN.CENTER)

    b._add_textbox(slide, Inches(5.6), Inches(2.9), Inches(2.5), Inches(0.3),
                   "← 温度波动时自动回退", Pt(11), RED_ALERT, align=PP_ALIGN.CENTER)

    b._add_textbox(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.35),
                   "状态转换规则", Pt(20), PRIMARY, bold=True)
    transitions = [
        "Idle -> Preparing：操作员点击[开始升温]，启动多阶段功率控制",
        "Preparing -> Ready：温度稳定（三重条件连续满足），系统自动转换",
        "Ready -> Preparing：温度波动超出稳定范围，自动回退重新升温",
        "Ready -> Recording：操作员点击[开始记录]，记录恒定功率并开始计时",
        "Recording -> Complete：用户手动停止 或 自动终止（30-60分钟检查点温漂<=0.5°C/10min）",
    ]
    for i, t in enumerate(transitions):
        b._add_textbox(slide, Inches(0.8), Inches(3.85 + i * 0.5), Inches(11.5), Inches(0.4),
                       f"  {t}", Pt(12), TEXT_BODY)
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(7.5), Inches(6.1), Inches(5.2), Inches(0.8),
                         "（此处需要插上试验流程状态机图 —— 对应论文图5-1）")

    # ================================================================
    # SLIDE 16: 3.3 数据库设计 (对应模板Slide 16)
    # ================================================================
    slide = b._content_slide("3. 系统总体设计 · 数据库设计", "Chapter 3 · 系统总体设计")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
                   "SQLite 6 张核心表", Pt(20), PRIMARY, bold=True)
    tables = [
        ["表名", "用途", "关键字段"],
        ["operators", "操作员账号", "userid, username, pwd, usertype"],
        ["apparatus", "试验设备", "apparatusid, name, pidport, constpower"],
        ["productmaster", "产品信息", "productid, productname, specific, diameter"],
        ["testmaster", "试验主记录", "testid, productid, ambtemp, totaltesttime"],
        ["sensors", "传感器配置", "sensorid, sensorname, signaltype, outputvalue"],
        ["CalibrationRecords", "校准记录", "CalibrationType, TemperatureData(JSON)"],
    ]
    b._add_table(slide, Inches(0.8), Inches(2.15),
                 [Inches(2.2), Inches(2.0), Inches(4.5)], tables[0], tables[1:], row_h=Inches(0.48))

    b._card(slide, Inches(0.8), Inches(5.5), Inches(3.7), Inches(1.4),
            "🌱", "种子数据", [
                "admin / 123456（管理员）",
                "experimenter / 123456（试验员）",
            ])
    b._card(slide, Inches(4.8), Inches(5.5), Inches(3.7), Inches(1.4),
            "🔗", "ER关系", [
                "Testmaster ↔ Productmaster 一对多",
                "ProductId 外键关联",
            ])
    b._card(slide, Inches(8.8), Inches(5.5), Inches(3.9), Inches(1.4),
            "✨", "设计亮点", [
                "配置驱动数据库路径",
                "UI层零SQL代码",
            ])
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(9.5), Inches(2.15), Inches(3.2), Inches(3.0),
                         "（此处需要插上数据库ER图 —— 对应论文图4-2）")

    # ================================================================
    # SLIDE 17: 3.4 功能模块设计 (对应模板Slide 17)
    # ================================================================
    slide = b._content_slide("3. 系统总体设计 · 功能模块", "Chapter 3 · 系统总体设计")
    modules = [
        ("🧠", "试验管理控制", "状态机驱动全流程\n自动升温→稳定判断→\n数据记录→报告生成", SECONDARY),
        ("🌡", "智能温度控制", "多阶段自适应算法\n+ PID闭环控制\n750°C ± 0.5°C 精度", ACCENT),
        ("📡", "数据采集与处理", "800ms高频采集\n双缓冲区策略\n内存缓存+数据库", ACCENT2),
        ("🖥", "用户交互", "实时曲线可视化\n状态驱动UI\n角色权限管理", PRIMARY),
        ("📄", "报告自动生成", "Excel格式化报表\nPDF正式报告\nCSV原始数据", RGBColor(0x8E, 0x44, 0xAD)),
        ("📐", "校准系统", "炉壁9点温度校验\n中心轴15点扫描\n校准前置管理", RGBColor(0xE7, 0x4C, 0x3C)),
    ]
    for i, (icon, title, desc, color) in enumerate(modules):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.7 + row * 2.3)
        b._card(slide, left, top, Inches(3.9), Inches(2.0), icon, title, [desc], color)
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.75),
                         "（此处需要插上系统功能结构图 —— 对应论文图4-3）")

    # ================================================================
    # SLIDE 18: Chapter 4 分隔页 — 系统详细设计与实现
    # ================================================================
    b._chapter_divider("04", "系统详细设计与实现", "SYSTEM IMPLEMENTATION")

    # ================================================================
    # SLIDE 19: 4.1 数据采集与通信 (对应模板Slide 19)
    # ================================================================
    slide = b._content_slide("4. 系统实现 · 数据采集服务", "Chapter 4 · 系统详细设计与实现")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "DaqWorker — 800ms 数据引擎", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "独立后台线程，800ms 固定周期轮询",
        "调用 TestController.Tick() 统一调度",
        "SensorSimulator 计算仿真温度数据",
        "AccumulateSensorData() 累积传感器缓冲",
        "EvaluateAutoTransitions() 状态自动迁移",
        "BuildSnapshot() 构建 RuntimeSnapshot 广播",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "数据流架构", Pt(20), PRIMARY, bold=True)
    flow_steps = [
        ("DaqWorker\n(800ms)", SECONDARY),
        ("TestController\n.Tick()", PRIMARY),
        ("SensorSimulator\n.Update()", SECONDARY),
        ("BuildSnapshot\n()", PRIMARY),
        ("DataBroadcast\n事件", ACCENT),
    ]
    for i, (label, color) in enumerate(flow_steps):
        left = Inches(7.0 + i * 1.25)
        b._flow_arrow(slide, left, Inches(2.3), Inches(1.1), label, color)
        if i < len(flow_steps) - 1:
            b._add_textbox(slide, left + Inches(1.12), Inches(2.38), Inches(0.15), Inches(0.4),
                           "→", Pt(14), ACCENT, bold=True)

    b._card(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.3),
            "🔌", "通信容错机制", [
                "FluentModbus 库封装 Modbus RTU 协议",
                "超时重试机制（默认1000ms）",
                "异常捕获与结构化日志记录",
                "断线自动重连 + 数据缓冲",
            ])
    b._card(slide, Inches(7.0), Inches(4.5), Inches(5.6), Inches(1.5),
            "📊", "5通道温度采集", [
                "TF1 — 炉壁热电偶1（蓝色曲线）",
                "TF2 — 炉壁热电偶2（红色曲线）",
                "TS/TC — 试样表面/中心热电偶",
            ])
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.75),
                         "（此处需要插上数据流设计图 —— 对应论文图5-4）")

    # ================================================================
    # SLIDE 20: 4.2 试验流程控制 (对应模板Slide 20)
    # ================================================================
    slide = b._content_slide("4. 系统实现 · 试验流程控制", "Chapter 4 · 系统详细设计与实现")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "TestController — 状态机控制器", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "TestState 枚举驱动 5 种状态",
        "Tick() 方法：800ms 统一调度入口",
        "EvaluateAutoTransitions()：自动判断迁移条件",
        "CheckAutoTermination()：自动终止判断",
        "Ready → Preparing 自动回退机制",
        "lock 保证线程安全",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "自动终止条件", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "30分钟检查点：温漂 ≤ 0.5°C/10min → 提前终止",
        "35分钟检查点：同上条件",
        "40/45/50/55分钟检查点：同上条件",
        "60分钟：无条件终止",
        "温漂由 MathNet.Numerics 线性回归计算",
    ], Pt(13))

    b._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "🧠", "恒定功率计算", [
                "Ready 状态记录 PID 输出值",
                "最多缓存 600 个采样点",
                "StartRecording 时取平均值",
                "作为 Recording 阶段恒定功率",
            ])
    b._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "📋", "传感器数据缓冲", [
                "SensorDataBuffer 累积时序数据",
                "每 Tick 追加 SensorDataRecord",
                "含12通道值 + 时间戳",
                "供导出模块读取完整试验数据",
            ])
    b._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(2.3),
            "🔄", "状态转换日志", [
                "每次转换生成 SystemMessage",
                "含时间戳 + 转换原因",
                "通过广播推送至 UI",
                "底部日志区实时显示",
            ], ACCENT2)

    # ================================================================
    # SLIDE 21: 4.3 智能温控算法实现 (对应模板Slide 21)
    # ================================================================
    slide = b._content_slide("4. 系统实现 · 智能温控算法", "Chapter 4 · 系统详细设计与实现")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
                   "SensorSimulator — 仿真引擎", Pt(20), PRIMARY, bold=True)

    phases = [
        ("升温阶段\n(Preparing)", "线性递增\nstep = HeatingRate × 0.8\nTF1/TF2 同步升温", SECONDARY),
        ("稳定阶段\n(Ready)", "钳位目标温度\n+ 随机噪声\n模拟真实波动", ACCENT2),
        ("记录阶段\n(Recording)", "炉温恒定\n表面温/中心温\n指数逼近炉温", ACCENT),
        ("冷却阶段\n(Cooling)", "指数衰减\n自然散热模拟\n回落至安全温度", RGBColor(0x95, 0xA5, 0xA6)),
    ]
    for i, (title, desc, color) in enumerate(phases):
        left = Inches(0.5 + i * 3.2)
        b._add_rounded_rect(slide, left, Inches(2.2), Inches(2.9), Inches(2.0), color)
        b._add_textbox(slide, left + Inches(0.15), Inches(2.3), Inches(2.6), Inches(0.6),
                       title, Pt(13), WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.2)
        b._add_textbox(slide, left + Inches(0.15), Inches(2.95), Inches(2.6), Inches(1.0),
                       desc, Pt(11), RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER, line_spacing=1.4)

    b._add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.35),
                   "温漂计算：MathNet.Numerics 线性回归", Pt(18), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(4.95), Inches(5.8), Inches(0.5), [
        "最近20个采样点用于线性回归",
        "Fit.Line(times, temps) → (intercept, slope)",
        "slope 即为温漂速率（°C/s）",
        "乘以600得到10分钟漂移量",
    ], Pt(12))
    b._bullet_list(slide, Inches(7.0), Inches(4.95), Inches(5.8), Inches(0.5), [
        "稳定性判断：温度范围 + 持续稳定计数",
        "IsTemperatureStable()：±3°C 阈值判断",
        "连续4次稳定 → 触发 Ready 状态",
        "仿真参数从 appsettings.json 读取",
    ], Pt(12))
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(7.0), Inches(6.1), Inches(5.8), Inches(0.75),
                         "（此处需要插上温控方案性能对比表 —— 对应论文表5-2）")

    # ================================================================
    # SLIDE 22: 4.4 客户端实现 (对应模板Slide 22)
    # ================================================================
    slide = b._content_slide("4. 系统实现 · 客户端界面", "Chapter 4 · 系统详细设计与实现")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
                   "WinForms 客户端核心功能", Pt(20), PRIMARY, bold=True)

    ui_features = [
        ("📊", "实时监控", "OxyPlot 温度曲线\n10分钟滑动窗口\n800ms 刷新周期\n鼠标悬停查坐标"),
        ("🎛", "流程控制", "状态驱动按钮矩阵\nButtonStateMatrix\n防误操作设计\n倒计时显示"),
        ("👤", "权限管理", "角色登录认证\nadmin/experimenter\nSHA256 密码哈希\n操作日志审计"),
        ("📐", "校准面板", "炉壁9点校验\n中心轴15点扫描\n校准前置管理\n历史记录查询"),
    ]
    for i, (icon, title, desc) in enumerate(ui_features):
        left = Inches(0.5 + i * 3.2)
        b._card(slide, left, Inches(2.2), Inches(2.9), Inches(2.5), icon, title, [desc])

    b._card(slide, Inches(0.8), Inches(5.0), Inches(3.7), Inches(1.0),
            "🔒", "线程安全设计", [
                "InvokeRequired + Invoke 安全调度",
            ])
    b._card(slide, Inches(4.8), Inches(5.0), Inches(3.7), Inches(1.0),
            "🎨", "界面设计亮点", [
                "OxyPlot 四通道实时温度曲线",
            ])
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(8.8), Inches(2.2), Inches(4.0), Inches(1.8),
                         "（此处需要插上试验监控主界面截图 —— 对应论文图5-5）")
    b._image_placeholder(slide, Inches(8.8), Inches(4.2), Inches(4.0), Inches(1.8),
                         "（此处需要插上新建试验界面截图 —— 对应论文图5-6）")

    # ================================================================
    # SLIDE 23: 4.5 导出与报表 (对应模板Slide 23)
    # ================================================================
    slide = b._content_slide("4. 系统实现 · 数据管理与导出", "Chapter 4 · 系统详细设计与实现")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "数据持久化流程", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "实时数据 → 内存缓冲 → 异步批量写入 SQLite",
        "CsvSampleWriter 同时输出 CSV 文件备份",
        "路径格式：TestData/{productid}/{testid}/sensor_data.csv",
        "TestRecordCoordinator 统一组织保存逻辑",
        "保存后 flag 置位，防止重复覆盖",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "多格式导出", Pt(20), PRIMARY, bold=True)
    exports = [
        ("📄", "CSV 导出", "原始传感器时序数据\n便于二次分析"),
        ("📊", "Excel 导出", "EPPlus 7.5.2\n格式化报表含图表"),
        ("📕", "PDF 导出", "PDFsharp-MigraDoc\n正式试验报告"),
    ]
    for i, (icon, title, desc) in enumerate(exports):
        left = Inches(7.0 + i * 2.1)
        b._card(slide, left, Inches(2.15), Inches(1.9), Inches(1.8), icon, title, [desc])

    b._add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.35),
                   "历史查询与数据管理", Pt(18), PRIMARY, bold=True)
    b._card(slide, Inches(0.8), Inches(4.75), Inches(5.6), Inches(2.0),
            "🔍", "HistoryCoordinator", [
                "按日期范围、产品编号、操作员多条件组合查询",
                "查看试验详细温度数据",
                "一键重新生成历史试验报告",
            ])
    b._card(slide, Inches(7.0), Inches(4.75), Inches(5.6), Inches(1.5),
            "📋", "试验记录管理", [
                "TestRecordDialog 录入火焰现象、质量损失等",
                "自动计算质量损失率、温升等指标",
            ])
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(0.8), Inches(6.4), Inches(5.6), Inches(0.75),
                         "（此处需要插上历史数据查询界面截图 —— 对应论文图5-9）")
    b._image_placeholder(slide, Inches(6.8), Inches(6.4), Inches(5.9), Inches(0.75),
                         "（此处需要插上记录查询/导出界面截图 —— 对应论文图5-10）")

    # ================================================================
    # SLIDE 24: 4.6 仿真模式实现 (对应模板Slide 24)
    # ================================================================
    slide = b._content_slide("4. 系统实现 · 仿真测试框架", "Chapter 4 · 系统详细设计与实现")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
                   "分层可配置仿真开关", Pt(20), PRIMARY, bold=True)

    modes = [
        ["模式", "EnableSimulation", "SimulateSensors", "SimulatePid", "说明"],
        ["全硬件", "false", "-", "-", "连接真实硬件设备"],
        ["全仿真", "true", "true", "true", "无需任何物理设备"],
        ["半实物A", "true", "true", "false", "仿真温度+真实PID"],
        ["半实物B", "true", "false", "true", "真实温度+仿真PID"],
    ]
    b._add_table(slide, Inches(0.8), Inches(2.15),
                 [Inches(1.3), Inches(2.0), Inches(2.0), Inches(1.8), Inches(3.0)],
                 modes[0], modes[1:], row_h=Inches(0.45))

    b._card(slide, Inches(0.8), Inches(4.8), Inches(5.6), Inches(2.0),
            "🔧", "配置驱动（appsettings.json）", [
                "Simulation 节点：仿真参数（目标温度、升温速率、稳定阈值）",
                "Hardware 节点：硬件参数（恒定功率、PID目标温度）",
                "Report 节点：报表输出路径和格式设置",
                "运行时修改参数无需重启应用",
            ])
    b._card(slide, Inches(7.0), Inches(4.8), Inches(5.6), Inches(1.4),
            "✅", "仿真测试价值", [
                "开发阶段：无需硬件即可完整测试核心业务逻辑",
                "培训阶段：安全环境供操作员学习使用",
            ], ACCENT2)
    # ---- 图片占位 ----
    b._image_placeholder(slide, Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.75),
                         "（此处需要插上仿真引擎温度曲线示意图 —— 对应论文图5-8）")

    # ================================================================
    # SLIDE 25: Chapter 5 — 系统测试与总结 (对应模板Slide 25)
    # ================================================================
    slide = b._content_slide("5. 系统测试与总结", "Chapter 5 · 系统测试与展望")

    # 测试金字塔
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.35),
                   "三层测试体系", Pt(20), PRIMARY, bold=True)
    pyramid = [
        ("UI 自动化测试\nFlaUI", "桌面自动化验收 · TC01-TC10",
         Inches(2.5), Inches(2.2), Inches(4.5), Inches(0.9), ACCENT),
        ("集成测试\nxUnit", "数据库 · CSV · 导出链路 · 历史查询",
         Inches(1.8), Inches(3.25), Inches(6.0), Inches(0.95), SECONDARY),
        ("单元测试\nxUnit", "状态机 · 仿真计算 · 按钮矩阵 · 认证 · 配置",
         Inches(1.0), Inches(4.35), Inches(8.0), Inches(1.0), PRIMARY),
    ]
    for label, desc, left, top, width, height, color in pyramid:
        b._add_rounded_rect(slide, left, top, width, height, color)
        b._add_textbox(slide, left + Inches(0.2), top + Inches(0.08), Inches(2.5), Inches(0.65),
                       label, Pt(12), WHITE, bold=True)
        b._add_textbox(slide, left + Inches(2.8), top + Inches(0.1), width - Inches(3.2),
                       Inches(0.65), desc, Pt(10), RGBColor(0xEE, 0xEE, 0xEE))

    # 总结
    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.35),
                   "项目成果与总结", Pt(20), PRIMARY, bold=True)
    results = [
        "✅ 全流程自动化：状态机驱动5状态精确流转",
        "✅ 智能温控：多阶段自适应 + PID闭环",
        "✅ 稳定性判断：线性拟合 + 三重条件",
        "✅ 实时可视化：OxyPlot 4通道温度曲线",
        "✅ 多格式导出：CSV/Excel/PDF 一键输出",
        "✅ 仿真框架：分层开关，脱机测试",
    ]
    for i, r in enumerate(results):
        b._add_textbox(slide, Inches(7.0), Inches(2.15 + i * 0.4), Inches(5.5), Inches(0.35),
                       r, Pt(12), TEXT_DARK)

    # 展望
    b._add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.3), PRIMARY)
    b._add_textbox(slide, Inches(1.2), Inches(5.7), Inches(10.5), Inches(0.35),
                   "未来展望", Pt(18), ACCENT, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(6.1), Inches(10.5), Inches(0.6),
                   "架构演进：向云端化、Web/移动端发展，实现多设备集中远程运维  |  "
                   "智能升级：引入机器学习预测材料不燃性等级  |  "
                   "平台扩展：模块化架构兼容 ISO 1716 等更多测试标准",
                   Pt(11), WHITE, line_spacing=1.5)
    # ---- 图片占位（放在测试金字塔与总结之间的右侧空白） ----
    b._image_placeholder(slide, Inches(9.5), Inches(5.6), Inches(3.3), Inches(1.3),
                         "（此处需要插上温度控制精度统计表 —— 对应论文表6-6）")

    # ================================================================
    # SLIDE 26: 致谢 (对应模板Slide 26)
    # ================================================================
    slide = b.prs.slides.add_slide(b.prs.slide_layouts[6])
    b._set_bg(slide, PRIMARY)
    b._add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    b._add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x08, 0x1F, 0x3A))
    b._add_textbox(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.2),
                   "敬请各位老师批评指正", Pt(48), WHITE, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(3.3), Inches(11), Inches(0.6),
                   "Thank You for Your Guidance", Pt(22), ACCENT, font_name=EN_FONT)
    b._add_rect(slide, Inches(1.2), Inches(4.1), Inches(3), Inches(0.04), ACCENT)
    b._add_textbox(slide, Inches(1.2), Inches(4.4), Inches(11), Inches(0.5),
                   "建筑材料不燃性实验自动化测控系统研究与实现", Pt(18), RGBColor(0xAA, 0xBB, 0xCC))
    b._add_textbox(slide, Inches(1.2), Inches(5.2), Inches(11), Inches(0.8),
                   "答辩人：罗凌岚  |  学号：202231060826  |  指导教师：肖斌教授\n"
                   "西南石油大学 · 计算机与软件学院 · 软件工程2022级",
                   Pt(14), RGBColor(0x88, 0x99, 0xAA), line_spacing=1.6)
    b.slide_num += 1

    # ================================================================
    # SAVE
    # ================================================================
    output_dir = r"D:\jineng\jinengshijain\jinengshijain\docs"
    output_path = os.path.join(output_dir, "ISO11820-答辩PPT-项目版.pptx")
    b.prs.save(output_path)
    print(f"PPT saved to: {output_path}")
    print(f"Total slides: {b.slide_num}")


if __name__ == "__main__":
    build()
