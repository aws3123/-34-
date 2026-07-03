#!/usr/bin/env python3
"""ISO 11820 答辩PPT — 鎏金色主题版本"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 鎏金色主题配色
# ============================================================
DARK_BG      = RGBColor(0x0C, 0x0C, 0x1E)   # 深邃墨蓝黑
DARK_CARD    = RGBColor(0x15, 0x15, 0x2A)   # 深色卡片
DARK_CARD2   = RGBColor(0x1E, 0x1E, 0x38)   # 稍亮深色
GOLD         = RGBColor(0xC9, 0xA8, 0x4C)   # 鎏金色（主色）
GOLD_LIGHT   = RGBColor(0xE8, 0xD5, 0xA3)   # 香槟金
GOLD_BRIGHT  = RGBColor(0xF5, 0xE6, 0xC8)   # 暖白金
GOLD_DARK    = RGBColor(0x8B, 0x7D, 0x3C)   # 暗金
ROSE_GOLD    = RGBColor(0xB8, 0x7B, 0x5E)   # 玫瑰金
WARM_WHITE   = RGBColor(0xF8, 0xF4, 0xEB)   # 暖白（内容页底）
CARD_WHITE   = RGBColor(0xFF, 0xFB, 0xF0)   # 暖白卡片
TEXT_ON_DARK = RGBColor(0xE8, 0xE0, 0xD0)   # 深色上文字
TEXT_BODY    = RGBColor(0x3D, 0x35, 0x2A)   # 暖白底上正文
TEXT_DARK    = RGBColor(0x2A, 0x22, 0x18)   # 标题深色
SUBTLE_GOLD  = RGBColor(0xD4, 0xC5, 0x9A)   # 淡金辅助线
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT   = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)

BODY_FONT = 'Microsoft YaHei'
EN_FONT   = 'Arial'
SLIDE_W   = Inches(13.333)
SLIDE_H   = Inches(7.5)

# ============================================================
class GoldPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.slide_num = 0

    # ---- 基础 ----
    def _bg(self, slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _rect(self, slide, l, t, w, h, color, alpha=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _rrect(self, slide, l, t, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _tb(self, slide, l, t, w, h, text, sz=Pt(18),
            color=TEXT_BODY, bold=False, align=PP_ALIGN.LEFT, fn=BODY_FONT,
            ls=1.3, af=False):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tx.text_frame.word_wrap = True
        if af: tx.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = tx.text_frame.paragraphs[0]
        p.text = text; p.font.size = sz; p.font.color.rgb = color
        p.font.bold = bold; p.font.name = fn; p.alignment = align
        p.line_spacing = ls; p.space_before = Pt(0); p.space_after = Pt(0)
        return tx

    def _mtb(self, slide, l, t, w, h, lines, dsz=Pt(16), dc=TEXT_BODY,
             db=False, ls=1.3, af=True):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tx.text_frame.word_wrap = True
        if af: tx.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf = tx.text_frame
        for i, ln in enumerate(lines):
            if isinstance(ln, str):
                txt, sz, c, b = ln, dsz, dc, db
            else:
                txt = ln[0]; sz = ln[1] if len(ln)>1 else dsz
                c = ln[2] if len(ln)>2 else dc; b = ln[3] if len(ln)>3 else db
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = txt; p.font.size = sz; p.font.color.rgb = c
            p.font.bold = b; p.font.name = BODY_FONT
            p.line_spacing = ls; p.space_after = Pt(3); p.space_before = Pt(0)
        return tx

    def _pn(self, slide):
        self.slide_num += 1
        self._tb(slide, Inches(12.0), Inches(7.05), Inches(1), Inches(0.35),
                 str(self.slide_num), Pt(11), GOLD_DARK, align=PP_ALIGN.RIGHT)

    # ---- 鎏金标题栏 ----
    def _gold_bar(self, slide, title, sub=None):
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.35), DARK_BG)
        # 底部金色细线
        self._rect(slide, Inches(0), Inches(1.35), SLIDE_W, Inches(0.035), GOLD)
        # 左侧金色装饰竖条
        self._rect(slide, Inches(0.5), Inches(0.25), Inches(0.06), Inches(0.85), GOLD)
        self._tb(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.85),
                 title, Pt(30), GOLD_LIGHT, bold=True)
        if sub:
            self._tb(slide, Inches(0.8), Inches(0.88), Inches(11), Inches(0.4),
                     sub, Pt(13), GOLD_DARK)

    # ---- 内容页（暖白底） ----
    def _content(self, title, sub=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, WARM_WHITE)
        self._gold_bar(slide, title, sub)
        # 右下角金色装饰
        self._rect(slide, Inches(12.8), Inches(6.8), Inches(0.5), Inches(0.035), GOLD)
        self._pn(slide)
        return slide

    # ---- 章节分隔页（深色底 + 大号金色数字） ----
    def _chapter(self, num, cn, en):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._bg(slide, DARK_BG)
        # 左侧金色竖线
        self._rect(slide, Inches(0), Inches(0), Inches(0.05), SLIDE_H, GOLD)
        # 顶部装饰金线
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), GOLD)
        # 底部装饰
        self._rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.04), GOLD_DARK)
        # 大号金色数字
        self._tb(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(2.2),
                 num, Pt(110), GOLD, bold=True, fn=EN_FONT)
        # 金色横线
        self._rect(slide, Inches(0.8), Inches(3.7), Inches(4), Inches(0.04), GOLD)
        # 中文标题
        self._tb(slide, Inches(0.8), Inches(4.0), Inches(9), Inches(0.9),
                 cn, Pt(38), GOLD_LIGHT, bold=True)
        # 英文
        self._tb(slide, Inches(0.8), Inches(4.9), Inches(9), Inches(0.5),
                 en, Pt(15), GOLD_DARK, fn=EN_FONT)
        # 右侧竖线
        self._rect(slide, Inches(12.5), Inches(0.8), Inches(0.025), Inches(5.8), GOLD_DARK)
        self._pn(slide)
        return slide

    # ---- 鎏金卡片 ----
    def _card(self, slide, l, t, w, h, icon, title, body, color=GOLD):
        # 深色卡片（在暖白底上）
        self._rrect(slide, l, t, w, h, CARD_WHITE)
        # 顶部金色装饰条
        self._rect(slide, l, t, w, Inches(0.05), color)
        # 左侧金色小竖条
        self._rect(slide, l, t, Inches(0.05), h, color)
        # icon
        self._tb(slide, l+Inches(0.2), t+Inches(0.12), Inches(0.5), Inches(0.4),
                 icon, Pt(20), color, bold=True)
        # title
        self._tb(slide, l+Inches(0.2), t+Inches(0.5), w-Inches(0.4), Inches(0.35),
                 title, Pt(14), TEXT_DARK, bold=True)
        # body
        self._mtb(slide, l+Inches(0.2), t+Inches(0.9), w-Inches(0.4), h-Inches(1.0),
                  body, dsz=Pt(11), dc=TEXT_BODY, ls=1.35)

    def _dark_card(self, slide, l, t, w, h, icon, title, body, color=GOLD):
        self._rrect(slide, l, t, w, h, DARK_CARD)
        self._rect(slide, l, t, w, Inches(0.05), color)
        self._tb(slide, l+Inches(0.2), t+Inches(0.12), Inches(0.5), Inches(0.4),
                 icon, Pt(20), color, bold=True)
        self._tb(slide, l+Inches(0.2), t+Inches(0.5), w-Inches(0.4), Inches(0.35),
                 title, Pt(14), GOLD_LIGHT, bold=True)
        self._mtb(slide, l+Inches(0.2), t+Inches(0.9), w-Inches(0.4), h-Inches(1.0),
                  body, dsz=Pt(11), dc=TEXT_ON_DARK, ls=1.35)

    def _bullets(self, slide, l, t, w, h, items, sz=Pt(14), c=TEXT_BODY, sp=1.4):
        lines = [f"  {x}" if isinstance(x, str) else f"  {x[0]}" for x in items]
        self._mtb(slide, l, t, w, h, lines, dsz=sz, dc=c, ls=sp)

    def _table(self, slide, l, t, cws, hds, rows, hc=DARK_BG, rh=Inches(0.48)):
        nr = len(rows)+1; nc = len(hds); tw = sum(cws)
        ts = slide.shapes.add_table(nr, nc, l, t, tw, rh*nr)
        tbl = ts.table
        for ci, cw in enumerate(cws): tbl.columns[ci].width = cw
        for ci, h in enumerate(hds):
            cell = tbl.cell(0, ci); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = hc
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=GOLD_LIGHT
                p.font.name=BODY_FONT; p.alignment=PP_ALIGN.CENTER
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci); cell.text = str(val)
                bg = RGBColor(0xF5,0xF0,0xE2) if ri%2==0 else CARD_WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size=Pt(11); p.font.color.rgb=TEXT_BODY
                    p.font.name=BODY_FONT; p.alignment=PP_ALIGN.CENTER
        return ts

    def _placeholder(self, slide, l, t, w, h, hint):
        PH_BG = RGBColor(0xFF, 0xF3, 0xCD)
        PH_BD = RED_ACCENT; PH_TX = RGBColor(0xC0, 0x39, 0x2B)
        s = self._rrect(slide, l, t, w, h, PH_BG)
        s.line.color.rgb = PH_BD; s.line.width = Pt(2.5)
        self._tb(slide, l, t+Inches(0.08), w, Inches(0.35),
                 "📷", Pt(16), PH_BD, bold=True, align=PP_ALIGN.CENTER)
        self._tb(slide, l+Inches(0.1), t+Inches(0.4), w-Inches(0.2), h-Inches(0.5),
                 hint, Pt(11), PH_TX, bold=True, align=PP_ALIGN.CENTER, ls=1.4)

    def _flow(self, slide, l, t, w, text, color=GOLD):
        s = self._rrect(slide, l, t, w, Inches(0.6), color)
        self._tb(slide, l, t+Inches(0.04), w, Inches(0.52),
                 text, Pt(10), WHITE if color!=GOLD else DARK_BG, bold=True,
                 align=PP_ALIGN.CENTER, ls=1.1)

    def _num_item(self, slide, l, t, w, n, title, desc, color=GOLD):
        badge = self._rrect(slide, l, t, Inches(0.42), Inches(0.42), color)
        self._tb(slide, l, t+Inches(0.02), Inches(0.42), Inches(0.38),
                 str(n), Pt(18), DARK_BG, bold=True, align=PP_ALIGN.CENTER)
        self._tb(slide, l+Inches(0.55), t, w-Inches(0.55), Inches(0.3),
                 title, Pt(15), TEXT_DARK, bold=True)
        self._tb(slide, l+Inches(0.55), t+Inches(0.32), w-Inches(0.55), Inches(0.5),
                 desc, Pt(11), TEXT_BODY, af=True)


# ============================================================
def build():
    g = GoldPPT()

    # ============ SLIDE 1: 封面 ============
    slide = g.prs.slides.add_slide(g.prs.slide_layouts[6])
    g._bg(slide, DARK_BG)
    # 顶部金线
    g._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), GOLD)
    # 底部金线
    g._rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.04), GOLD_DARK)
    g._rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), RGBColor(0x06,0x06,0x12))
    # 左侧金色装饰
    g._rect(slide, Inches(0.8), Inches(1.6), Inches(0.06), Inches(2.8), GOLD)
    # 主标题
    g._tb(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(1.2),
          "建筑材料不燃性实验", Pt(46), GOLD_LIGHT, bold=True)
    g._tb(slide, Inches(1.2), Inches(2.7), Inches(11), Inches(1.0),
          "自动化测控系统研究与实现", Pt(42), GOLD, bold=True)
    # 副标题
    g._tb(slide, Inches(1.2), Inches(3.9), Inches(11), Inches(0.5),
          "—— 毕业论文答辩", Pt(22), GOLD_DARK)
    # 金色分隔线
    g._rect(slide, Inches(1.2), Inches(4.6), Inches(3.5), Inches(0.03), GOLD)
    # 技术标签
    g._tb(slide, Inches(1.2), Inches(4.9), Inches(11), Inches(0.4),
          ".NET 8  ·  WinForms  ·  SQLite  ·  OxyPlot  ·  EPPlus  ·  MathNet.Numerics",
          Pt(13), GOLD_DARK)
    # 作者信息
    g._tb(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
          "答辩人：罗凌岚  |  指导教师：肖斌教授  |  2026年5月",
          Pt(14), TEXT_ON_DARK)
    g._tb(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.4),
          "西南石油大学 · 计算机与软件学院 · 软件工程2022级",
          Pt(12), GOLD_DARK)
    # 右侧装饰金色圆角矩形
    g._rrect(slide, Inches(10.5), Inches(1.8), Inches(2.3), Inches(3.5), DARK_CARD)
    g._rect(slide, Inches(10.5), Inches(1.8), Inches(2.3), Inches(0.04), GOLD)
    g._tb(slide, Inches(10.7), Inches(2.1), Inches(1.9), Inches(0.5),
          "ISO", Pt(28), GOLD, bold=True, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g._tb(slide, Inches(10.7), Inches(2.6), Inches(1.9), Inches(0.5),
          "11820", Pt(36), GOLD_LIGHT, bold=True, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g._rect(slide, Inches(10.9), Inches(3.2), Inches(1.5), Inches(0.02), GOLD_DARK)
    g._tb(slide, Inches(10.7), Inches(3.4), Inches(1.9), Inches(0.8),
          "2020", Pt(24), GOLD_DARK, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g._tb(slide, Inches(10.7), Inches(4.1), Inches(1.9), Inches(0.5),
          "Non-combustibility", Pt(9), GOLD_DARK, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g.slide_num += 1

    # ============ SLIDE 2: 目录 ============
    slide = g._content("汇报大纲", "CONTENTS")
    agenda = [
        ("01", "研究背景与意义", "研究背景、国内外现状、研究目标"),
        ("02", "相关技术与理论", "温控算法、稳定性判断、开发技术"),
        ("03", "系统总体设计", "架构设计、状态机、数据库、功能模块"),
        ("04", "系统详细设计与实现", "服务端实现、客户端实现、仿真模式"),
        ("05", "系统测试与总结", "功能测试、实验数据、总结与展望"),
    ]
    for i, (n, t, d) in enumerate(agenda):
        col = i % 3; row = i // 3
        l = Inches(0.6 + col * 4.2); t_ = Inches(2.0 + row * 2.2)
        g._num_item(slide, l, t_, Inches(3.8), n, t, d)

    # ============ SLIDE 3: Ch1 分隔 ============
    g._chapter("01", "研究背景与意义", "BACKGROUND AND SIGNIFICANCE")

    # ============ SLIDE 4: 研究背景 ============
    slide = g._content("1. 研究背景与意义 · 研究背景", "Chapter 1")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
          "研究背景", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "建筑材料防火安全性能是构筑社会公共安全防线的基石",
        "不燃性是评价建筑材料防火等级最基础、最重要的性能指标",
        "ISO 11820-2020 标准规定：750C 高温条件下连续稳定10分钟",
        "传统人工测试模式效率低、过程可控性不强、数据一致性不佳",
    ], Pt(13))
    g._tb(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
          "传统人工测试的痛点", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "操作人员需全程密切关注并手动调节，劳动强度极大",
        "凭经验判断温度稳定，难以实现高精度控制",
        "人工记录和计算数据效率低，容易引入主观误差",
        "无法满足 ISO 11820-2020 对过程控制精度的严格条件",
    ], Pt(13))
    # 研究目标卡
    g._rrect(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.2), DARK_BG)
    g._rect(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.04), GOLD)
    g._tb(slide, Inches(1.1), Inches(4.8), Inches(2), Inches(0.35),
          "研究目标", Pt(18), GOLD, bold=True)
    g._tb(slide, Inches(1.1), Inches(5.3), Inches(10.8), Inches(1.2),
          "开发一套严格依照 ISO 11820-2020 国际标准、能替代人工操作达成全流程自动化控制的测控系统，"
          "以大幅提升建筑材料不燃性测试的效率、准确性以及数据规范性。",
          Pt(14), TEXT_ON_DARK, ls=1.5)

    # ============ SLIDE 5: 国内外现状 ============
    slide = g._content("1. 研究背景与意义 · 国内外现状", "Chapter 1")
    g._card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0),
            "CN", "国内研究现状", [
                "GB/T 5464-2010 标准为实验搭建了基础框架",
                "自动化程度较低，主要依赖人工操作",
                "有学者探索测试设备改进：自动定位传感器",
                "模块化设计提升了试验稳定性",
                "但在 C# 高级语言开发的智能测控方面案例较少",
                "难以满足实时数据分析和远程监控的需求",
            ])
    g._card(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(5.0),
            "EU", "国外研究现状", [
                "欧美主导，聚焦智能化与自动化技术深度整合",
                "数据驱动方法突显自动化系统在提升精度的必要性",
                "AI/ML/DL 技术在建筑4.0中广泛应用",
                "NIST 提出融合红外热成像的动态监测平台",
                "但缺乏针对 ISO 11820 标准的定制化方案",
                "面临高成本和数据兼容性挑战",
            ], ROSE_GOLD)

    # ============ SLIDE 6: 研究意义 ============
    slide = g._content("1. 研究背景与意义 · 研究意义", "Chapter 1")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
          "本研究的核心价值", Pt(20), GOLD, bold=True)
    vals = [
        ("T", "技术层面", "为类似标准检测方法的自动化改造\n提供可复用的技术框架和工程范例", GOLD),
        ("D", "数据层面", "标准化、结构化数据为实验室信息化\n管理与材料性能大数据分析奠定基础", ROSE_GOLD),
        ("E", "效率层面", "试验效率提升3-5倍，人工干预时间\n从3-4小时减少到约10分钟", GOLD_DARK),
        ("I", "行业层面", "推动检测行业从经验判断朝\n数据驱动转型升级", GOLD_LIGHT),
    ]
    for i, (ic, t, d, c) in enumerate(vals):
        g._card(slide, Inches(0.5+i*3.15), Inches(2.3), Inches(2.9), Inches(2.5), ic, t, [d], c)
    g._rrect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), DARK_BG)
    g._rect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.04), GOLD)
    g._tb(slide, Inches(1.2), Inches(5.35), Inches(10.5), Inches(0.35),
          "研究方法", Pt(18), GOLD, bold=True)
    g._tb(slide, Inches(1.2), Inches(5.8), Inches(10.5), Inches(0.7),
          "运用工业测控技术和现代软件工程方法，基于 .NET 8 框架与 WinForms 三层架构，"
          "集成状态机模式、多阶段自适应温控算法和线性拟合稳定性判断算法，"
          "设计并实现一套针对建筑材料不燃性试验的自动化测控系统。",
          Pt(13), TEXT_ON_DARK, ls=1.5)

    # ============ SLIDE 7: Ch2 分隔 ============
    g._chapter("02", "相关技术与理论", "KEY TECHNOLOGIES AND THEORIES")

    # ============ SLIDE 8: 多阶段温控算法 ============
    slide = g._content("2. 关键技术 · 多阶段自适应温控算法", "Chapter 2")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
          "算法核心思想", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "采用分段逼近、平滑过渡的控制策略",
        "升温过程划分为五个逻辑阶段",
        "各阶段运用不同的功率输出策略",
        "融合开环响应快与闭环精度高的优势",
    ], Pt(13))
    stages = [
        ["温度区间", "功率策略", "控制方式"],
        ["室温~300C", "30% 额定功率(软启动)", "开环"],
        ["300~500C", "50% 额定功率", "开环"],
        ["500~600C", "70% 额定功率", "开环"],
        ["600~700C", "90% 额定功率", "开环"],
        [">700C", "切换 PID 闭环", "闭环"],
    ]
    g._table(slide, Inches(7.0), Inches(1.7),
             [Inches(1.8), Inches(2.2), Inches(1.5)], stages[0], stages[1:])
    g._card(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(1.5),
            "K", "模式切换", [
                "当炉温突破 700C 阈值后执行关键模式切换",
                "通过 Modbus 协议向 PID 控制器写入目标温度",
                "将控制权移交至硬件 PID 闭环控制器",
            ])
    g._card(slide, Inches(7.0), Inches(4.5), Inches(5.6), Inches(1.5),
            "V", "算法优势", [
                "最大超调量仅 2.1C（全功率方案为 32.6C）",
                "稳态波动范围 +-0.5C 以内",
                "升温速度与超调抑制之间取得最优平衡",
            ], GREEN_ACCENT)
    g._placeholder(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.75),
                   "（此处需要插上多阶段升温过程功率阶梯示意图 —— 对应论文图2-1）")

    # ============ SLIDE 9: 稳定性判断 ============
    slide = g._content("2. 关键技术 · 温度稳定性自动判断", "Chapter 2")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
          "数学基础：最小二乘线性回归", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "动态滑动窗口：存储最近10分钟温度数据（600采样点）",
        "采样周期 800ms，每分钟调用 CheckStartCriteria()",
        "使用 MathNet.Numerics 库的 Fit.Line() 函数",
        "计算拟合直线斜率，得到温度漂移量",
    ], Pt(13))
    g._tb(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
          "三重条件联合判定", Pt(20), GOLD, bold=True)
    conds = [["判据", "条件", "ISO标准要求"],
             ["温度范围", "T in [745, 755]", "目标温度 +-5C"],
             ["温度漂移", "Drift<=2.0C/10min", "<=2C/10min"],
             ["温度偏差", "Tmax-Tavg<=1.0C", "<=1C"]]
    g._table(slide, Inches(7.0), Inches(2.15),
             [Inches(1.5), Inches(2.3), Inches(1.8)], conds[0], conds[1:])
    g._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(1.5),
            "L", "持续性验证机制", [
                "每项条件维护独立倒计时计数器(初始600)",
                "条件满足则计数器减1，否则重置为600",
                "三个计数器同时归零才判定稳定",
            ])
    g._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(1.5),
            "M", "线性拟合 vs 其他方法", [
                "极差法：对单次野值极其敏感",
                "方差法：无法识别缓慢单调漂移",
                "线性拟合：抗干扰性 + 趋势敏感性",
            ])
    g._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(1.5),
            "S", "核心优势", [
                "将ISO标准定性要求转化为精确数学模型",
                "全自动、客观、连续验证",
            ], GREEN_ACCENT)
    g._placeholder(slide, Inches(0.8), Inches(6.15), Inches(5.6), Inches(0.75),
                   "（此处需要插上线性拟合算法核心代码截图 —— 对应论文图5-2）")
    g._placeholder(slide, Inches(6.8), Inches(6.15), Inches(5.9), Inches(0.75),
                   "（此处需要插上三重条件判断核心代码截图 —— 对应论文图5-3）")

    # ============ SLIDE 10: 技术栈 ============
    slide = g._content("2. 关键技术 · 开发技术栈", "Chapter 2")
    techs = [
        (".N", ".NET 8 / C# 10", ["跨平台高性能框架", "WinForms 桌面客户端", "异步编程模型"]),
        ("OP", "OxyPlot 2.1.2", ["实时温度曲线绘制", "10分钟滑动窗口", "自适应数据抽稀"]),
        ("SQ", "SQLite + EF Core", ["轻量级嵌入式数据库", "6张核心业务表", "ORM 对象关系映射"]),
        ("EP", "EPPlus + PDFsharp", ["Excel 格式化报表", "PDF 正式报告", "CSV 原始数据"]),
        ("MN", "MathNet.Numerics", ["线性回归分析", "温度漂移计算", "数值计算库"]),
        ("SR", "Serilog 4.0", ["结构化日志记录", "文件滚动存储", "故障追踪审计"]),
    ]
    for i, (ic, t, ls) in enumerate(techs):
        col = i%3; row = i//3
        g._card(slide, Inches(0.5+col*4.2), Inches(1.7+row*2.8),
                Inches(3.9), Inches(2.5), ic, t, ls)

    # ============ SLIDE 11: 通信与设计模式 ============
    slide = g._content("2. 关键技术 · 通信协议与设计模式", "Chapter 2")
    g._card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(2.5),
            "MB", "Modbus RTU 通信协议", [
                "工业标准串行通信协议，广泛应用于温控设备",
                "功能码 0x03：读取寄存器（温度数据）",
                "功能码 0x06：写入寄存器（目标温度、控制模式）",
                "FluentModbus 库封装，超时重试 + 异常处理",
            ])
    g._card(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(2.5),
            "SM", "状态机设计模式", [
                "将复杂试验流程离散为有限个状态",
                "状态间转换由明确定义的事件触发",
                "确保试验严格按 ISO 11820 标准步骤执行",
                "代码结构清晰，可维护性好，易于扩展",
            ], ROSE_GOLD)
    g._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "AR", "三层架构设计", [
                "表示层：WinForms UI 界面",
                "业务逻辑层：状态机 + 算法",
                "数据访问层：SQLite + EF Core",
                "高内聚、低耦合",
            ])
    g._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "TS", "线程安全机制", [
                "后台线程驱动数据采集（800ms）",
                "MainForm 通过 InvokeRequired 检测",
                "Invoke 安全调度 UI 更新",
                "RuntimeSnapshot 不可变快照",
            ])
    g._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(2.3),
            "SIM", "仿真测试支持", [
                "虚拟串口 + Modbus Slave 模拟器",
                "无需物理硬件即可完整测试",
                "分层可配置仿真开关",
                "极大降低开发验证成本",
            ], GREEN_ACCENT)

    # ============ SLIDE 12: 技术总结 ============
    slide = g._content("2. 关键技术 · 核心技术体系", "Chapter 2")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
          "核心技术体系总览", Pt(20), GOLD, bold=True)
    cores = [
        ("多阶段自适应\n温控算法", GOLD, "室温->750C\n五阶段功率控制\n+ PID 闭环切换"),
        ("线性拟合\n稳定性判断", ROSE_GOLD, "MathNet.Numerics\n三重条件联合判定\n600秒持续验证"),
        ("状态机\n流程控制", GOLD_DARK, "5状态精确建模\n自动迁移+回退\nISO标准严格执行"),
        ("Modbus RTU\n硬件通信", GOLD_LIGHT, "PID温控器通信\n传感器数据采集\n超时重试容错"),
        ("数据持久化\n与报表", RGBColor(0x8E,0x44,0xAD), "SQLite 6张表\nCSV/Excel/PDF\n自动化报告生成"),
    ]
    for i, (t, c, d) in enumerate(cores):
        l = Inches(0.3+i*2.6)
        g._rrect(slide, l, Inches(2.3), Inches(2.4), Inches(1.6), c)
        g._tb(slide, l+Inches(0.15), Inches(2.4), Inches(2.1), Inches(0.65),
              t, Pt(13), DARK_BG, bold=True, align=PP_ALIGN.CENTER, ls=1.2)
        g._tb(slide, l+Inches(0.15), Inches(3.1), Inches(2.1), Inches(0.7),
              d, Pt(10), DARK_BG, align=PP_ALIGN.CENTER, ls=1.3)
    g._tb(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.35),
          "技术选型理由", Pt(18), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(4.75), Inches(5.8), Inches(0.5), [
        ".NET 8：统一平台，高性能运行时，跨平台支持",
        "WinForms：成熟桌面GUI框架，控件丰富",
        "SQLite：轻量嵌入式，无需独立服务",
        "OxyPlot：高性能图表库，支持实时数据流",
    ], Pt(12))
    g._bullets(slide, Inches(7.0), Inches(4.75), Inches(5.8), Inches(0.5), [
        "MathNet.Numerics：权威数值计算库",
        "EPPlus/PDFsharp：成熟Office文档生成方案",
        "Serilog：结构化日志，多目标输出",
        "FluentModbus：简洁的Modbus协议封装",
    ], Pt(12))

    # ============ SLIDE 13: Ch3 分隔 ============
    g._chapter("03", "系统总体设计", "SYSTEM ARCHITECTURE DESIGN")

    # ============ SLIDE 14: 系统架构 ============
    slide = g._content("3. 系统总体设计 · 系统架构", "Chapter 3")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.35),
          "四层技术架构", Pt(20), GOLD, bold=True)
    layers = [
        ("表示层 (UI)", "WinForms 图形界面 · 温度曲线 · 状态显示", GOLD),
        ("业务逻辑层", "TestController 状态机 · 温控算法 · 稳定性判断", ROSE_GOLD),
        ("数据访问层", "SQLite + EF Core · 数据持久化 · 导出服务", RGBColor(0x8E,0x44,0xAD)),
        ("硬件抽象层", "Modbus RTU · PID控制器 · 传感器采集 · 仿真引擎", GOLD_DARK),
    ]
    for i, (n, d, c) in enumerate(layers):
        top = Inches(2.2+i*1.0)
        g._rrect(slide, Inches(0.8), top, Inches(6.5), Inches(0.82), c)
        g._tb(slide, Inches(1.1), top+Inches(0.08), Inches(3), Inches(0.35),
              n, Pt(16), DARK_BG, bold=True)
        g._tb(slide, Inches(4.0), top+Inches(0.1), Inches(3.2), Inches(0.35),
              d, Pt(11), DARK_BG)
    g._tb(slide, Inches(0.8), Inches(6.3), Inches(7.5), Inches(0.35),
          "Core 零外部依赖 | UI 不写业务逻辑 | 共享契约跨层交互",
          Pt(12), GOLD, bold=True, align=PP_ALIGN.CENTER)
    g._placeholder(slide, Inches(7.8), Inches(2.0), Inches(5.0), Inches(4.5),
                   "（此处需要插上系统四层架构图 —— 对应论文图4-1）")

    # ============ SLIDE 15: 状态机 ============
    slide = g._content("3. 系统总体设计 · 状态机设计", "Chapter 3")
    states = [("Idle","空闲",RGBColor(0x95,0xA5,0xA6)),
              ("Preparing","升温中",GOLD_DARK),
              ("Ready","就绪",GREEN_ACCENT),
              ("Recording","记录中",ROSE_GOLD),
              ("Complete","完成",DARK_BG)]
    for i, (code, label, color) in enumerate(states):
        l = Inches(0.5+i*2.55)
        g._rrect(slide, l, Inches(1.9), Inches(2.2), Inches(0.95), color)
        g._tb(slide, l, Inches(1.98), Inches(2.2), Inches(0.42),
              code, Pt(18), GOLD_LIGHT if color==DARK_BG else WHITE,
              bold=True, align=PP_ALIGN.CENTER)
        g._tb(slide, l, Inches(2.4), Inches(2.2), Inches(0.35),
              label, Pt(12), GOLD_DARK if color==DARK_BG else RGBColor(0xEE,0xEE,0xEE),
              align=PP_ALIGN.CENTER)
        if i < len(states)-1:
            g._tb(slide, l+Inches(2.25), Inches(2.1), Inches(0.3), Inches(0.4),
                  ">", Pt(26), GOLD, bold=True, align=PP_ALIGN.CENTER)
    g._tb(slide, Inches(5.6), Inches(2.9), Inches(2.5), Inches(0.3),
          "<-- 温度波动时自动回退", Pt(11), RED_ACCENT, align=PP_ALIGN.CENTER)
    g._tb(slide, Inches(0.8), Inches(3.4), Inches(6.5), Inches(0.35),
          "状态转换规则", Pt(20), GOLD, bold=True)
    trans = [
        "Idle -> Preparing：操作员点击[开始升温]，启动多阶段功率控制",
        "Preparing -> Ready：温度稳定（三重条件连续满足），系统自动转换",
        "Ready -> Preparing：温度波动超出稳定范围，自动回退重新升温",
        "Ready -> Recording：操作员点击[开始记录]，记录恒定功率并开始计时",
        "Recording -> Complete：用户手动停止 或 自动终止（30-60分钟检查点）",
    ]
    for i, t in enumerate(trans):
        g._tb(slide, Inches(0.8), Inches(3.85+i*0.48), Inches(6.5), Inches(0.4),
              t, Pt(11), TEXT_BODY)
    g._placeholder(slide, Inches(7.8), Inches(3.4), Inches(5.0), Inches(3.2),
                   "（此处需要插上试验流程状态机图 —— 对应论文图5-1）")

    # ============ SLIDE 16: 数据库设计 ============
    slide = g._content("3. 系统总体设计 · 数据库设计", "Chapter 3")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.35),
          "SQLite 6 张核心表", Pt(20), GOLD, bold=True)
    tbls = [["表名", "用途", "关键字段"],
            ["operators", "操作员账号", "id, username, pwd, role"],
            ["apparatus", "试验设备", "id, name, model, serial_number"],
            ["productmaster", "产品信息", "id, product_code, product_name"],
            ["testmaster", "试验主记录", "productid+testid(PK), testdate"],
            ["sensors", "传感器配置", "id, name, type, channel"],
            ["CalibrationRecords", "校准记录", "id, sensor_id(FK), result_json"]]
    g._table(slide, Inches(0.8), Inches(2.15),
             [Inches(2.0), Inches(1.6), Inches(3.5)], tbls[0], tbls[1:])
    g._card(slide, Inches(0.8), Inches(5.5), Inches(3.7), Inches(1.4),
            "SD", "种子数据", ["admin / 123456（管理员）", "experimenter / 123456（试验员）"])
    g._card(slide, Inches(4.8), Inches(5.5), Inches(3.7), Inches(1.4),
            "ER", "ER关系", ["Testmaster <-> Productmaster 一对多", "ProductId 外键关联"])
    g._card(slide, Inches(8.8), Inches(5.5), Inches(3.9), Inches(1.4),
            "HL", "设计亮点", ["配置驱动数据库路径", "UI层零SQL代码"])
    # ER图占位（如果之前已插入则跳过，这里保留指示）
    g._placeholder(slide, Inches(9.5), Inches(2.15), Inches(3.2), Inches(3.0),
                   "（此处需要插上数据库ER图 —— 对应论文图4-2）")

    # ============ SLIDE 17: 功能模块 ============
    slide = g._content("3. 系统总体设计 · 功能模块", "Chapter 3")
    mods = [
        ("TQ", "试验管理控制", "状态机驱动全流程\n自动升温->稳定判断->\n数据记录->报告生成", GOLD),
        ("TC", "智能温度控制", "多阶段自适应算法\n+ PID闭环控制\n750C +-0.5C 精度", ROSE_GOLD),
        ("DQ", "数据采集与处理", "800ms高频采集\n双缓冲区策略\n内存缓存+数据库", GREEN_ACCENT),
        ("UI", "用户交互", "实时曲线可视化\n状态驱动UI\n角色权限管理", GOLD_DARK),
        ("RP", "报告自动生成", "Excel格式化报表\nPDF正式报告\nCSV原始数据", RGBColor(0x8E,0x44,0xAD)),
        ("CL", "校准系统", "炉壁9点温度校验\n中心轴15点扫描\n校准前置管理", RED_ACCENT),
    ]
    for i, (ic, t, d, c) in enumerate(mods):
        col = i%3; row = i//3
        g._card(slide, Inches(0.5+col*4.2), Inches(1.7+row*2.3),
                Inches(3.9), Inches(2.0), ic, t, [d], c)
    g._placeholder(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.75),
                   "（此处需要插上系统功能结构图 —— 对应论文图4-3）")

    # ============ SLIDE 18: Ch4 分隔 ============
    g._chapter("04", "系统详细设计与实现", "SYSTEM IMPLEMENTATION")

    # ============ SLIDE 19: 数据采集 ============
    slide = g._content("4. 系统实现 · 数据采集服务", "Chapter 4")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
          "DaqWorker -- 800ms 数据引擎", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "独立后台线程，800ms 固定周期轮询",
        "调用 TestController.Tick() 统一调度",
        "SensorSimulator 计算仿真温度数据",
        "AccumulateSensorData() 累积传感器缓冲",
        "EvaluateAutoTransitions() 状态自动迁移",
        "BuildSnapshot() 构建 RuntimeSnapshot 广播",
    ], Pt(13))
    g._tb(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
          "数据流架构", Pt(20), GOLD, bold=True)
    flows = [("DaqWorker\n(800ms)",GOLD_DARK), ("TestController\n.Tick()",DARK_BG),
             ("SensorSim\n.Update()",GOLD), ("BuildSnap\n()",DARK_BG), ("DataBroadcast\n事件",ROSE_GOLD)]
    for i, (lb, c) in enumerate(flows):
        l = Inches(7.0+i*1.25)
        g._flow(slide, l, Inches(2.3), Inches(1.1), lb, c)
        if i < len(flows)-1:
            g._tb(slide, l+Inches(1.12), Inches(2.38), Inches(0.15), Inches(0.4),
                  ">", Pt(14), GOLD, bold=True)
    g._card(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(1.5),
            "CM", "通信容错机制", [
                "FluentModbus 库封装 Modbus RTU 协议",
                "超时重试机制（默认1000ms）",
                "异常捕获与结构化日志记录",
            ])
    g._card(slide, Inches(7.0), Inches(4.5), Inches(5.6), Inches(1.5),
            "5C", "5通道温度采集", [
                "TF1 — 炉壁热电偶1（蓝色曲线）",
                "TF2 — 炉壁热电偶2（红色曲线）",
                "TS/TC — 试样表面/中心热电偶",
            ])
    g._placeholder(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.75),
                   "（此处需要插上数据流设计图 —— 对应论文图5-4）")

    # ============ SLIDE 20: 试验流程控制 ============
    slide = g._content("4. 系统实现 · 试验流程控制", "Chapter 4")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
          "TestController -- 状态机控制器", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "TestState 枚举驱动 5 种状态",
        "Tick() 方法：800ms 统一调度入口",
        "EvaluateAutoTransitions()：自动判断迁移条件",
        "CheckAutoTermination()：自动终止判断",
        "Ready -> Preparing 自动回退机制",
        "lock 保证线程安全",
    ], Pt(13))
    g._tb(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
          "自动终止条件", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "30分钟检查点：温漂<=0.5C/10min -> 提前终止",
        "35/40/45/50/55分钟检查点：同上条件",
        "60分钟：无条件终止",
        "温漂由 MathNet.Numerics 线性回归计算",
    ], Pt(13))
    g._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "CP", "恒定功率计算", [
                "Ready 状态记录 PID 输出值",
                "最多缓存 600 个采样点",
                "StartRecording 时取平均值",
                "作为 Recording 阶段恒定功率",
            ])
    g._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "SB", "传感器数据缓冲", [
                "SensorDataBuffer 累积时序数据",
                "每 Tick 追加 SensorDataRecord",
                "含12通道值 + 时间戳",
                "供导出模块读取完整试验数据",
            ])
    g._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(2.3),
            "LOG", "状态转换日志", [
                "每次转换生成 SystemMessage",
                "含时间戳 + 转换原因",
                "通过广播推送至 UI",
                "底部日志区实时显示",
            ], GREEN_ACCENT)

    # ============ SLIDE 21: 温控算法实现 ============
    slide = g._content("4. 系统实现 · 智能温控算法", "Chapter 4")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
          "SensorSimulator -- 仿真引擎", Pt(20), GOLD, bold=True)
    phases = [
        ("升温阶段\n(Preparing)", "线性递增\nstep=Rate*0.8\nTF1/TF2同步", GOLD_DARK),
        ("稳定阶段\n(Ready)", "钳位目标温度\n+随机噪声\n模拟真实波动", GREEN_ACCENT),
        ("记录阶段\n(Recording)", "炉温恒定\n表面温/中心温\n指数逼近炉温", ROSE_GOLD),
        ("冷却阶段\n(Cooling)", "指数衰减\n自然散热模拟\n回落至安全温度", RGBColor(0x95,0xA5,0xA6)),
    ]
    for i, (t, d, c) in enumerate(phases):
        l = Inches(0.5+i*3.2)
        g._rrect(slide, l, Inches(2.2), Inches(2.9), Inches(2.0), c)
        g._tb(slide, l+Inches(0.15), Inches(2.3), Inches(2.6), Inches(0.6),
              t, Pt(13), WHITE, bold=True, align=PP_ALIGN.CENTER, ls=1.2)
        g._tb(slide, l+Inches(0.15), Inches(2.95), Inches(2.6), Inches(1.0),
              d, Pt(11), RGBColor(0xEE,0xEE,0xEE), align=PP_ALIGN.CENTER, ls=1.4)
    g._tb(slide, Inches(0.8), Inches(4.5), Inches(5.8), Inches(0.35),
          "温漂计算：MathNet.Numerics 线性回归", Pt(18), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(4.95), Inches(5.8), Inches(0.5), [
        "最近20个采样点用于线性回归",
        "Fit.Line(times, temps) -> (intercept, slope)",
        "slope 即为温漂速率（C/s）",
        "乘以600得到10分钟漂移量",
    ], Pt(12))
    g._bullets(slide, Inches(7.0), Inches(4.95), Inches(5.8), Inches(0.5), [
        "稳定性判断：温度范围 + 持续稳定计数",
        "IsTemperatureStable()：+-3C 阈值判断",
        "连续4次稳定 -> 触发 Ready 状态",
        "仿真参数从 appsettings.json 读取",
    ], Pt(12))
    g._placeholder(slide, Inches(7.0), Inches(6.1), Inches(5.8), Inches(0.75),
                   "（此处需要插上温控方案性能对比表 —— 对应论文表5-2）")

    # ============ SLIDE 22: 客户端界面 ============
    slide = g._content("4. 系统实现 · 客户端界面", "Chapter 4")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(8), Inches(0.35),
          "WinForms 客户端核心功能", Pt(20), GOLD, bold=True)
    uifs = [
        ("RT", "实时监控", "OxyPlot温度曲线\n10分钟滑动窗口\n800ms刷新周期"),
        ("FC", "流程控制", "状态驱动按钮矩阵\nButtonStateMatrix\n防误操作设计"),
        ("AU", "权限管理", "角色登录认证\nadmin/experimenter\nSHA256密码哈希"),
        ("CL", "校准面板", "炉壁9点校验\n中心轴15点扫描\n校准前置管理"),
    ]
    for i, (ic, t, d) in enumerate(uifs):
        g._card(slide, Inches(0.5+i*2.1), Inches(2.2), Inches(1.9), Inches(2.2), ic, t, [d])
    g._card(slide, Inches(0.8), Inches(4.7), Inches(3.7), Inches(1.0),
            "TS", "线程安全设计", ["InvokeRequired + Invoke 安全调度"])
    g._card(slide, Inches(4.8), Inches(4.7), Inches(3.7), Inches(1.0),
            "HL", "界面设计亮点", ["OxyPlot 四通道实时温度曲线"])
    g._placeholder(slide, Inches(8.8), Inches(2.2), Inches(4.0), Inches(1.8),
                   "（此处需要插上试验监控主界面截图 —— 对应论文图5-5）")
    g._placeholder(slide, Inches(8.8), Inches(4.2), Inches(4.0), Inches(1.8),
                   "（此处需要插上新建试验界面截图 —— 对应论文图5-6）")

    # ============ SLIDE 23: 数据管理与导出 ============
    slide = g._content("4. 系统实现 · 数据管理与导出", "Chapter 4")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
          "数据持久化流程", Pt(20), GOLD, bold=True)
    g._bullets(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "实时数据 -> 内存缓冲 -> 异步批量写入 SQLite",
        "CsvSampleWriter 同时输出 CSV 文件备份",
        "路径格式：TestData/{productid}/{testid}/sensor_data.csv",
        "TestRecordCoordinator 统一组织保存逻辑",
        "保存后 flag 置位，防止重复覆盖",
    ], Pt(13))
    g._tb(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
          "多格式导出", Pt(20), GOLD, bold=True)
    exps = [("CSV", "CSV 导出", "原始传感器时序数据"),
            ("XLS", "Excel 导出", "EPPlus 格式化报表"),
            ("PDF", "PDF 导出", "PDFsharp正式报告")]
    for i, (ic, t, d) in enumerate(exps):
        g._card(slide, Inches(7.0+i*2.1), Inches(2.15), Inches(1.9), Inches(1.5), ic, t, [d])
    g._card(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(1.5),
            "HQ", "HistoryCoordinator", [
                "按日期范围、产品编号、操作员多条件查询",
                "查看试验详细温度数据，一键重新生成报告",
            ])
    g._card(slide, Inches(7.0), Inches(4.5), Inches(5.6), Inches(1.5),
            "TR", "试验记录管理", [
                "TestRecordDialog 录入火焰现象、质量损失等",
                "自动计算质量损失率、温升等指标",
            ])
    g._placeholder(slide, Inches(0.8), Inches(6.15), Inches(5.6), Inches(0.75),
                   "（此处需要插上历史数据查询界面截图 —— 对应论文图5-9）")
    g._placeholder(slide, Inches(6.8), Inches(6.15), Inches(5.9), Inches(0.75),
                   "（此处需要插上记录查询/导出界面截图 —— 对应论文图5-10）")

    # ============ SLIDE 24: 仿真模式 ============
    slide = g._content("4. 系统实现 · 仿真测试框架", "Chapter 4")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.35),
          "分层可配置仿真开关", Pt(20), GOLD, bold=True)
    modes = [["模式", "EnableSim", "SimSensors", "SimPid", "说明"],
             ["全硬件", "false", "-", "-", "连接真实硬件设备"],
             ["全仿真", "true", "true", "true", "无需任何物理设备"],
             ["半实物A", "true", "true", "false", "仿真温度+真实PID"],
             ["半实物B", "true", "false", "true", "真实温度+仿真PID"]]
    g._table(slide, Inches(0.8), Inches(2.15),
             [Inches(1.3), Inches(1.8), Inches(1.8), Inches(1.5), Inches(3.0)],
             modes[0], modes[1:])
    g._card(slide, Inches(0.8), Inches(4.8), Inches(5.6), Inches(1.4),
            "CFG", "配置驱动（appsettings.json）", [
                "Simulation 节点：仿真参数（目标温度、升温速率）",
                "Hardware 节点：硬件参数（恒定功率、PID温度）",
                "运行时修改参数无需重启应用",
            ])
    g._card(slide, Inches(7.0), Inches(4.8), Inches(5.6), Inches(1.4),
            "VL", "仿真测试价值", [
                "开发阶段：无需硬件即可完整测试核心业务逻辑",
                "培训阶段：安全环境供操作员学习使用",
            ], GREEN_ACCENT)
    g._placeholder(slide, Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.75),
                   "（此处需要插上仿真引擎温度曲线示意图 —— 对应论文图5-8）")

    # ============ SLIDE 25: 测试与总结 ============
    slide = g._content("5. 系统测试与总结", "Chapter 5")
    g._tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.35),
          "三层测试体系", Pt(20), GOLD, bold=True)
    pyr = [
        ("UI 自动化\nFlaUI", "桌面自动化验收 TC01-TC10",
         Inches(2.5), Inches(2.2), Inches(4.5), Inches(0.9), ROSE_GOLD),
        ("集成测试\nxUnit", "数据库 CSV 导出链路 历史查询",
         Inches(1.8), Inches(3.25), Inches(6.0), Inches(0.95), GOLD),
        ("单元测试\nxUnit", "状态机 仿真计算 按钮矩阵 认证 配置",
         Inches(1.0), Inches(4.35), Inches(8.0), Inches(1.0), DARK_BG),
    ]
    for lb, desc, l, t, w, h, c in pyr:
        g._rrect(slide, l, t, w, h, c)
        g._tb(slide, l+Inches(0.2), t+Inches(0.08), Inches(2.5), Inches(0.65),
              lb, Pt(12), GOLD_LIGHT if c==DARK_BG else WHITE, bold=True)
        g._tb(slide, l+Inches(2.8), t+Inches(0.1), w-Inches(3.2), Inches(0.65),
              desc, Pt(10), GOLD_DARK if c==DARK_BG else RGBColor(0xEE,0xEE,0xEE))
    g._tb(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.35),
          "项目成果与总结", Pt(20), GOLD, bold=True)
    results = [
        ">> 全流程自动化：状态机驱动5状态精确流转",
        ">> 智能温控：多阶段自适应 + PID闭环",
        ">> 稳定性判断：线性拟合 + 三重条件",
        ">> 实时可视化：OxyPlot 4通道温度曲线",
        ">> 多格式导出：CSV/Excel/PDF 一键输出",
        ">> 仿真框架：分层开关，脱机测试",
    ]
    for i, r in enumerate(results):
        g._tb(slide, Inches(7.0), Inches(2.15+i*0.4), Inches(5.5), Inches(0.35),
              r, Pt(12), TEXT_DARK)
    g._rrect(slide, Inches(0.8), Inches(5.6), Inches(8.5), Inches(1.3), DARK_BG)
    g._rect(slide, Inches(0.8), Inches(5.6), Inches(8.5), Inches(0.04), GOLD)
    g._tb(slide, Inches(1.2), Inches(5.7), Inches(7.5), Inches(0.35),
          "未来展望", Pt(18), GOLD, bold=True)
    g._tb(slide, Inches(1.2), Inches(6.1), Inches(7.5), Inches(0.6),
          "架构演进：向云端化、Web/移动端发展 | "
          "智能升级：引入机器学习预测材料不燃性等级 | "
          "平台扩展：模块化架构兼容 ISO 1716 等更多测试标准",
          Pt(11), TEXT_ON_DARK, ls=1.5)
    g._placeholder(slide, Inches(9.5), Inches(5.6), Inches(3.3), Inches(1.3),
                   "（此处需要插上温度控制精度统计表 —— 对应论文表6-6）")

    # ============ SLIDE 26: 致谢 ============
    slide = g.prs.slides.add_slide(g.prs.slide_layouts[6])
    g._bg(slide, DARK_BG)
    g._rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), GOLD)
    g._rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.04), GOLD_DARK)
    g._rect(slide, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), RGBColor(0x06,0x06,0x12))
    # 左侧金色竖线装饰
    g._rect(slide, Inches(1.0), Inches(1.8), Inches(0.06), Inches(3.0), GOLD)
    g._tb(slide, Inches(1.4), Inches(2.0), Inches(10), Inches(1.0),
          "敬请各位老师批评指正", Pt(48), GOLD_LIGHT, bold=True)
    g._tb(slide, Inches(1.4), Inches(3.2), Inches(10), Inches(0.5),
          "Thank You for Your Guidance", Pt(22), GOLD_DARK, fn=EN_FONT)
    g._rect(slide, Inches(1.4), Inches(4.0), Inches(3), Inches(0.03), GOLD)
    g._tb(slide, Inches(1.4), Inches(4.3), Inches(10), Inches(0.5),
          "建筑材料不燃性实验自动化测控系统研究与实现", Pt(18), GOLD_DARK)
    g._tb(slide, Inches(1.4), Inches(5.0), Inches(10), Inches(0.8),
          "答辩人：罗凌岚  |  学号：202231060826  |  指导教师：肖斌教授\n"
          "西南石油大学 · 计算机与软件学院 · 软件工程2022级",
          Pt(14), TEXT_ON_DARK, ls=1.6)
    # 右侧 ISO 标识
    g._rrect(slide, Inches(10.5), Inches(2.0), Inches(2.2), Inches(3.0), DARK_CARD)
    g._rect(slide, Inches(10.5), Inches(2.0), Inches(2.2), Inches(0.04), GOLD)
    g._tb(slide, Inches(10.6), Inches(2.3), Inches(2.0), Inches(0.5),
          "ISO", Pt(26), GOLD, bold=True, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g._tb(slide, Inches(10.6), Inches(2.8), Inches(2.0), Inches(0.5),
          "11820", Pt(32), GOLD_LIGHT, bold=True, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g._rect(slide, Inches(10.8), Inches(3.4), Inches(1.6), Inches(0.02), GOLD_DARK)
    g._tb(slide, Inches(10.6), Inches(3.6), Inches(2.0), Inches(0.5),
          "2020", Pt(22), GOLD_DARK, fn=EN_FONT, align=PP_ALIGN.CENTER)
    g.slide_num += 1

    # ============ SAVE ============
    out = os.path.join(r"D:\jineng\jinengshijain\jinengshijain\docs",
                       "ISO11820-答辩PPT-鎏金版.pptx")
    g.prs.save(out)
    print(f"PPT saved: {out}")
    print(f"Total slides: {g.slide_num}")


if __name__ == "__main__":
    build()
