#!/usr/bin/env python3
"""ISO 11820 课程项目答辩 PPT 生成器 — 修复文本溢出"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============================================================
# 配色方案 - 科技深蓝
# ============================================================
PRIMARY = RGBColor(0x0D, 0x2B, 0x4E)
SECONDARY = RGBColor(0x1A, 0x56, 0x9E)
ACCENT = RGBColor(0xE8, 0x8D, 0x2A)
ACCENT2 = RGBColor(0x27, 0xAE, 0x60)
LIGHT_BG = RGBColor(0xF0, 0xF3, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
TEXT_BODY = RGBColor(0x34, 0x49, 0x5E)
SUBTLE_LINE = RGBColor(0xBD, 0xC3, 0xD0)
CARD_BG = RGBColor(0xFA, 0xFB, 0xFD)

BODY_FONT = 'Microsoft YaHei'
EN_FONT = 'Arial'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class PPTBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.slide_num = 0

    def _set_bg(self, slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _add_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_rounded_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_textbox(self, slide, left, top, width, height, text, font_size=Pt(18),
                     color=TEXT_BODY, bold=False, align=PP_ALIGN.LEFT, font_name=BODY_FONT,
                     line_spacing=1.3, auto_fit=False):
        """Add a textbox. If auto_fit=True, height is minimum and box expands down."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        if auto_fit:
            txBox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = txBox.text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        return txBox

    def _add_multiline_textbox(self, slide, left, top, width, height, lines,
                               default_size=Pt(18), default_color=TEXT_BODY,
                               default_bold=False, line_spacing=1.3, auto_fit=True):
        """lines: list of str or (text, size, color, bold). auto_fit expands height."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        if auto_fit:
            txBox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf = txBox.text_frame
        for i, line in enumerate(lines):
            if isinstance(line, str):
                text, size, color, bold = line, default_size, default_color, default_bold
            else:
                text = line[0]
                size = line[1] if len(line) > 1 else default_size
                color = line[2] if len(line) > 2 else default_color
                bold = line[3] if len(line) > 3 else default_bold
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = size
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = BODY_FONT
            p.line_spacing = line_spacing
            p.space_after = Pt(3)
            p.space_before = Pt(0)
        return txBox

    def _page_number(self, slide):
        self.slide_num += 1
        self._add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35),
                          str(self.slide_num), Pt(11), SUBTLE_LINE, align=PP_ALIGN.RIGHT)

    def _title_bar(self, slide, title, subtitle=None):
        self._add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.35), PRIMARY)
        self._add_rect(slide, Inches(0), Inches(1.35), SLIDE_W, Inches(0.04), ACCENT)
        self._add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.85),
                          title, Pt(30), WHITE, bold=True)
        if subtitle:
            self._add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.4),
                              subtitle, Pt(14), RGBColor(0xBB, 0xCC, 0xDD))

    def _content_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, LIGHT_BG)
        self._title_bar(slide, title, subtitle)
        self._page_number(slide)
        return slide

    def _card(self, slide, left, top, width, height, icon_text, title, body_lines, color=SECONDARY):
        """Card with auto-fit body text. height is minimum."""
        self._add_rounded_rect(slide, left, top, width, height, WHITE)
        self._add_rect(slide, left, top, width, Inches(0.06), color)
        self._add_textbox(slide, left + Inches(0.25), top + Inches(0.15), Inches(0.5), Inches(0.45),
                          icon_text, Pt(22), color, bold=True)
        self._add_textbox(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.5),
                          Inches(0.35), title, Pt(14), TEXT_DARK, bold=True)
        self._add_multiline_textbox(slide, left + Inches(0.25), top + Inches(0.95),
                                    width - Inches(0.5), height - Inches(1.1),
                                    body_lines, default_size=Pt(11), default_color=TEXT_BODY,
                                    line_spacing=1.35, auto_fit=True)

    def _bullet_list(self, slide, left, top, width, height, items, size=Pt(15),
                     color=TEXT_BODY, spacing=1.4):
        lines = [f"▸ {item}" if isinstance(item, str) else f"▸ {item[0]}" for item in items]
        self._add_multiline_textbox(slide, left, top, width, height, lines,
                                    default_size=size, default_color=color,
                                    line_spacing=spacing, auto_fit=True)

    def _numbered_item(self, slide, left, top, width, num, title, desc, color=SECONDARY):
        badge = self._add_rounded_rect(slide, left, top, Inches(0.42), Inches(0.42), color)
        self._add_textbox(slide, left, top + Inches(0.02), Inches(0.42), Inches(0.38),
                          str(num), Pt(18), WHITE, bold=True, align=PP_ALIGN.CENTER)
        self._add_textbox(slide, left + Inches(0.55), top, width - Inches(0.55), Inches(0.3),
                          title, Pt(15), TEXT_DARK, bold=True)
        self._add_textbox(slide, left + Inches(0.55), top + Inches(0.32), width - Inches(0.55),
                          Inches(0.5), desc, Pt(11), TEXT_BODY, auto_fit=True)

    def _add_table(self, slide, left, top, col_widths, headers, rows, header_color=PRIMARY, row_h=Inches(0.5)):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        total_w = sum(col_widths)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w,
                                             row_h * n_rows)
        table = table_shape.table
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw
        for ci, h in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = BODY_FONT
                p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = str(val)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF1)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = TEXT_BODY
                    p.font.name = BODY_FONT
                    p.alignment = PP_ALIGN.CENTER
        return table_shape

    def _flow_arrow(self, slide, left, top, width, text, color=SECONDARY):
        shape = self._add_rounded_rect(slide, left, top, width, Inches(0.65), color)
        self._add_textbox(slide, left, top + Inches(0.04), width, Inches(0.57),
                          text, Pt(10), WHITE, bold=True, align=PP_ALIGN.CENTER,
                          line_spacing=1.1)

    def _arrow_right(self, slide, left, top):
        self._add_textbox(slide, left, top + Inches(0.08), Inches(0.25), Inches(0.45),
                          "→", Pt(18), ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
def build():
    b = PPTBuilder()

    # ================================================================
    # SLIDE 1: COVER
    # ================================================================
    slide = b.prs.slides.add_slide(b.prs.slide_layouts[6])
    b._set_bg(slide, PRIMARY)
    b._add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    b._add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x08, 0x1F, 0x3A))
    b._add_textbox(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.5),
                   "ISO 11820 材料燃烧试验仿真系统", Pt(44), WHITE, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(3.5), Inches(11), Inches(0.6),
                   "课程项目答辩", Pt(28), ACCENT, bold=False)
    b._add_rect(slide, Inches(1.2), Inches(4.3), Inches(2.5), Inches(0.04), ACCENT)
    b._add_textbox(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(0.5),
                   ".NET 8  WinForms  SQLite  OxyPlot  EPPlus  PDFsharp", Pt(14),
                   RGBColor(0x99, 0xAA, 0xBB))
    b._add_textbox(slide, Inches(1.2), Inches(5.5), Inches(5), Inches(0.4),
                   "集成负责人汇报  |  2026年6月", Pt(14), RGBColor(0x88, 0x99, 0xAA))
    b.slide_num += 1

    # ================================================================
    # SLIDE 2: AGENDA
    # ================================================================
    slide = b._content_slide("汇报大纲")
    agenda_items = [
        ("01", "项目概述", "背景、目标、应用场景、技术栈"),
        ("02", "系统架构总览", "分层设计、数据流、设计原则"),
        ("03", "核心业务流程", "状态机流转、全流程演示"),
        ("04", "运行时引擎", "状态机、PID仿真、广播机制"),
        ("05", "数据持久化与查询", "SQLite设计、历史查询、校准管理"),
        ("06", "导出与产物输出", "试验记录、CSV/Excel/PDF导出"),
        ("07", "团队协作与工程实践", "5人分工、边界规则、开发流程"),
        ("08", "测试体系", "单元/集成/UI自动化三层测试"),
        ("09", "成果展示与总结", "成果、收获、不足与展望"),
    ]
    for i, (num, title, desc) in enumerate(agenda_items):
        col = i % 3
        row = i // 3
        left = Inches(0.6 + col * 4.2)
        top = Inches(1.8 + row * 1.75)
        b._numbered_item(slide, left, top, Inches(3.8), num, title, desc)

    # ================================================================
    # SLIDE 3: CHAPTER 1 — 项目概述
    # ================================================================
    slide = b._content_slide("项目背景与定位", "Chapter 1 · 项目概述")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.35),
                   "项目定位", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(0.5), [
        "ISO 11820 材料燃烧试验仿真系统",
        "替代真实燃烧试验的桌面仿真工具",
        "模拟材料燃烧试验全流程：炉温控制、数据采集、试验记录、报告导出",
        "降低试验成本，保证数据可追溯",
    ], Pt(14))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "核心目标", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "真实模拟试验全流程",
        "5通道温度实时采集与显示",
        "自动状态机管理试验生命周期",
        "完整的试验记录与报告导出",
        "支持历史查询与设备校准",
    ], Pt(14))

    # ================================================================
    # SLIDE 4: CHAPTER 1 — 技术栈
    # ================================================================
    slide = b._content_slide("技术栈一览", "Chapter 1 · 项目概述")
    techs = [
        ("🖥", ".NET 8 / WinForms", ["桌面应用框架", "Windows Forms UI", "C# 12.0 语言"]),
        ("📊", "OxyPlot 2.1.2", ["温度曲线图表", "5通道实时绘制", "PlotView 控件"]),
        ("🗄", "SQLite", ["轻量级嵌入式数据库", "6张核心业务表", "Microsoft.Data.Sqlite"]),
        ("📄", "EPPlus + PDFsharp", ["Excel 报表导出", "PDF 报告生成", "CSV 数据输出"]),
        ("📝", "Serilog 4.0", ["结构化日志", "文件滚动存储", "10MB 自动分割"]),
        ("🧪", "xUnit + FlaUI", ["单元/集成测试", "桌面自动化验收", "UI Automation"]),
    ]
    for i, (icon, title, lines) in enumerate(techs):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.7 + row * 2.8)
        b._card(slide, left, top, Inches(3.9), Inches(2.5), icon, title, lines)

    # ================================================================
    # SLIDE 5: CHAPTER 2 — 分层架构
    # ================================================================
    slide = b._content_slide("分层架构设计", "Chapter 2 · 系统架构总览")
    # 设计原则
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.35),
                   "三大设计原则", Pt(20), PRIMARY, bold=True)
    principles = [
        "Core 不依赖 WinForms / SQLite / 文件系统",
        "UI 不直接写业务逻辑（无 SQL、无状态机判断）",
        "共享契约通过 Shared/Models、Shared/Events 交互",
    ]
    for i, p in enumerate(principles):
        b._add_textbox(slide, Inches(0.8 + i * 4.1), Inches(2.15), Inches(3.8), Inches(0.5),
                       f"0{i+1}  {p}", Pt(12), TEXT_BODY)

    # 分层架构
    layers = [
        ("ISO11820.Tests", "单元测试 / 集成测试", RGBColor(0x95, 0xA5, 0xA6)),
        ("ISO11820.UI.Tests", "FlaUI 桌面自动化测试", RGBColor(0x7F, 0x8C, 0x8D)),
        ("ISO11820.App", "WinForms 宿主 · 组合根 · Features · Infrastructure · Runtime", SECONDARY),
        ("ISO11820.Core", "纯业务领域 · 枚举 · 模型 · 契约（零外部依赖）", PRIMARY),
    ]
    for i, (name, desc, color) in enumerate(layers):
        top = Inches(2.9 + i * 1.05)
        b._add_rounded_rect(slide, Inches(0.8), top, Inches(11.5), Inches(0.85), color)
        b._add_textbox(slide, Inches(1.1), top + Inches(0.1), Inches(3.5), Inches(0.35),
                       name, Pt(17), WHITE, bold=True, font_name=EN_FONT)
        b._add_textbox(slide, Inches(4.8), top + Inches(0.12), Inches(7.2), Inches(0.35),
                       desc, Pt(13), RGBColor(0xEE, 0xEE, 0xEE))

    # ================================================================
    # SLIDE 6: CHAPTER 2 — 数据流
    # ================================================================
    slide = b._content_slide("核心数据流", "Chapter 2 · 系统架构总览")
    steps = [
        ("DaqWorker\n(800ms)", SECONDARY),
        ("TestController\n.Tick()", PRIMARY),
        ("SensorSimulator\n.Update()", SECONDARY),
        ("BuildSnapshot\n()", PRIMARY),
        ("DataBroadcast\n事件", ACCENT),
        ("MainForm\n.OnDataBroadcast", SECONDARY),
        ("UI 控件更新\n(Invoke安全)", PRIMARY),
    ]
    for i, (label, color) in enumerate(steps):
        left = Inches(0.25 + i * 1.86)
        b._flow_arrow(slide, left, Inches(2.5), Inches(1.7), label, color)
        if i < len(steps) - 1:
            b._arrow_right(slide, left + Inches(1.72), Inches(2.5))

    b._add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.35),
                   "数据流说明", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(3.95), Inches(5.8), Inches(0.5), [
        "DaqWorker 以 800ms 节拍驱动仿真引擎",
        "TestController.Tick() 统一调度状态迁移",
        "SensorSimulator 基于 PID 控制计算温度",
        "RuntimeSnapshot 聚合温度、状态、消息",
    ], Pt(13))
    b._bullet_list(slide, Inches(7.0), Inches(3.95), Inches(5.8), Inches(0.5), [
        "DataBroadcast 事件向所有订阅者广播",
        "MainForm 通过 InvokeRequired 回到 UI 线程",
        "确保跨线程安全，无 UI 冻结",
        "各模块仅消费快照，不自行计算状态",
    ], Pt(13))

    # ================================================================
    # SLIDE 7: CHAPTER 3 — 状态机
    # ================================================================
    slide = b._content_slide("状态机流转", "Chapter 3 · 核心业务流程")
    states = [
        ("Idle", "空闲", RGBColor(0x95, 0xA5, 0xA6)),
        ("Preparing", "准备中", SECONDARY),
        ("Ready", "就绪", ACCENT2),
        ("Recording", "记录中", ACCENT),
        ("Complete", "完成", PRIMARY),
    ]
    for i, (code, label, color) in enumerate(states):
        left = Inches(0.5 + i * 2.55)
        b._add_rounded_rect(slide, left, Inches(2.0), Inches(2.2), Inches(0.95), color)
        b._add_textbox(slide, left, Inches(2.08), Inches(2.2), Inches(0.42),
                       code, Pt(18), WHITE, bold=True, align=PP_ALIGN.CENTER)
        b._add_textbox(slide, left, Inches(2.5), Inches(2.2), Inches(0.35),
                       label, Pt(12), RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)
        if i < len(states) - 1:
            b._add_textbox(slide, left + Inches(2.25), Inches(2.2), Inches(0.3), Inches(0.4),
                           "→", Pt(26), ACCENT, bold=True, align=PP_ALIGN.CENTER)

    b._add_textbox(slide, Inches(5.6), Inches(3.0), Inches(2.5), Inches(0.3),
                   "← 温度不稳定时自动回退", Pt(11), RGBColor(0xE7, 0x4C, 0x3C), align=PP_ALIGN.CENTER)

    b._add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.35),
                   "状态说明", Pt(20), PRIMARY, bold=True)
    state_desc = [
        "Idle — 系统启动后初始状态，等待操作员新建试验",
        "Preparing — 试验已创建，开始升温，系统监控温度爬升",
        "Ready — 温度达到目标并稳定，等待操作员开始记录",
        "Recording — 正在记录试验数据，持续采集5通道温度",
        "Complete — 试验完成，可填写试验记录并导出报告",
    ]
    for i, desc in enumerate(state_desc):
        col = i % 3
        row = i // 3
        b._add_textbox(slide, Inches(0.8 + col * 4.1), Inches(4.05 + row * 0.55),
                       Inches(3.9), Inches(0.5), f"▸ {desc}", Pt(11), TEXT_BODY)

    # ================================================================
    # SLIDE 8: CHAPTER 3 — 全流程
    # ================================================================
    slide = b._content_slide("试验全流程", "Chapter 3 · 核心业务流程")
    flow = [
        ("1", "登录认证", "操作员/管理员\n角色登录"),
        ("2", "新建试验", "填写样品信息\n直径/高度/时长"),
        ("3", "开始升温", "5通道温度\n实时曲线绘制"),
        ("4", "稳定就绪", "温度达到目标\n系统自动判定"),
        ("5", "开始记录", "持续采集数据\n计算质量损失"),
        ("6", "停止加热", "结束试验\n进入冷却阶段"),
        ("7", "填写记录", "火焰现象/时长\n试验后质量/备注"),
        ("8", "导出报告", "CSV/Excel/PDF\n多格式导出"),
    ]
    for i, (num, title, desc) in enumerate(flow):
        left = Inches(0.25 + i * 1.62)
        b._numbered_item(slide, left, Inches(1.8), Inches(1.45), num, title, desc)

    b._add_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.03), SUBTLE_LINE)
    b._add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.35),
                   "关键数据：5通道温度  |  升温速率  |  稳定判据（温漂 < 2°C/min）  |  记录时长  |  质量损失率",
                   Pt(12), SECONDARY, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 9: CHAPTER 4 — 状态机 + 仿真
    # ================================================================
    slide = b._content_slide("状态机与仿真模型", "Chapter 4 · 运行时引擎")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "TestController — 唯一状态入口", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "TestState 枚举驱动 5 种状态",
        "EvaluateAutoTransitions() 自动判断迁移条件",
        "800ms Tick 驱动状态机推进",
        "状态迁移条件：升温→稳定(达到目标)、稳定→就绪(温漂<2°C/min)",
        "就绪→记录(操作员触发)、记录→完成(停止/自动终止)",
        "Ready → Preparing 自动回退机制",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "SensorSimulator — 仿真模型", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "800ms 节拍推进仿真数据",
        "PID 控制算法模拟升温过程",
        "线性回归计算温漂（斜率判定稳定）",
        "5通道独立温度仿真",
        "仿真参数从 appsettings.json 读取",
        "参数：目标温度、升温速率、稳定阈值、最大时长",
    ], Pt(13))

    # ================================================================
    # SLIDE 10: CHAPTER 4 — 广播 + 亮点
    # ================================================================
    slide = b._content_slide("广播机制与设计亮点", "Chapter 4 · 运行时引擎")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "广播链路", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "DaqWorker 定时器线程（800ms）",
        "→ TestController.Tick() 统一调度",
        "→ SensorSimulator.Update() 计算温度",
        "→ EvaluateAutoTransitions() 判断状态",
        "→ BuildSnapshot() 构建 RuntimeSnapshot",
        "→ DataBroadcast 事件广播到所有订阅者",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "线程安全机制", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "广播在后台线程触发",
        "MainForm 通过 InvokeRequired 检测",
        "Invoke 将更新调度到 UI 线程",
        "禁止后台线程直接操作控件",
        "RuntimeSnapshot 是不可变快照",
        "避免跨线程数据竞争",
    ], Pt(13))

    b._card(slide, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "🔌", "状态机与 UI 解耦", [
                "UI 只消费快照，不自行计算状态流转",
                "按钮状态由 ButtonStateMatrix 统一管理",
                "不拼装温度结构，不写状态判断",
            ])
    b._card(slide, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.3),
            "📡", "广播契约统一", [
                "RuntimeSnapshot 聚合所有运行时数据",
                "DataBroadcastEventArgs 承载快照",
                "各模块通过事件订阅，新增消费者无需改广播端",
            ])
    b._card(slide, Inches(8.8), Inches(4.5), Inches(3.9), Inches(2.3),
            "⚙", "参数可配置", [
                "appsettings.json 统一管理所有仿真参数",
                "目标温度、升温速率、稳定判据阈值",
                "不硬编码到仿真逻辑中，修改无需重新编译",
            ])

    # ================================================================
    # SLIDE 11: CHAPTER 5 — 数据库设计
    # ================================================================
    slide = b._content_slide("数据库设计与初始化", "Chapter 5 · 数据持久化与查询")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.35),
                   "SQLite 6 张核心表", Pt(20), PRIMARY, bold=True)
    tables = [
        ["表名", "用途", "关键字段"],
        ["operators", "操作员账号", "username, password, role"],
        ["apparatus", "试验设备", "apparatus_id, name, status"],
        ["productmaster", "样品主数据", "product_id, name, material"],
        ["testmaster", "试验主记录", "test_id, product_id, state, flag"],
        ["sensors", "传感器配置", "sensor_id, channel, type"],
        ["CalibrationRecords", "校准记录", "sensor_id, calib_data(JSON)"],
    ]
    b._add_table(slide, Inches(0.8), Inches(2.15),
                 [Inches(2.2), Inches(2.0), Inches(4.0)], tables[0], tables[1:], row_h=Inches(0.5))

    b._card(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.4),
            "🌱", "种子数据", [
                "admin / 123456（管理员）",
                "experimenter / 123456（试验员）",
                "默认设备记录 + 默认传感器记录",
            ])
    b._card(slide, Inches(7.0), Inches(5.5), Inches(5.5), Inches(1.4),
            "✨", "亮点", [
                "配置驱动数据库路径，不硬编码",
                "SQL 集中在 Infrastructure/Persistence，UI 层零 SQL",
                "DatabaseInitializer 可重复执行，不重复插入",
            ])

    # ================================================================
    # SLIDE 12: CHAPTER 5 — 历史查询与校准
    # ================================================================
    slide = b._content_slide("历史查询与校准管理", "Chapter 5 · 数据持久化与查询")
    b._card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0),
            "🔍", "HistoryCoordinator — 历史查询", [
                "按日期范围查询试验记录",
                "按样品编号 / 操作员筛选",
                "支持多条件组合查询",
                "查询结果：试验ID、样品、状态、时间、操作员",
                "与 TestRecordCoordinator 协作获取完整试验数据",
            ])
    b._card(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(5.0),
            "📐", "CalibrationCoordinator — 校准管理", [
                "多点校准记录读写",
                "CalibrationRecords 表存储校准参数",
                "JSON 字段灵活存储校准数据",
                "校准数据关联传感器，支持历史追溯",
                "校准面板实时温度显示，与广播协作获取当前温度",
                "校准数据用于传感器误差修正",
            ])

    # ================================================================
    # SLIDE 13: CHAPTER 6 — 导出
    # ================================================================
    slide = b._content_slide("试验记录保存与导出", "Chapter 6 · 导出与产物输出")
    b._card(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(2.5),
            "📝", "TestRecordCoordinator — 记录保存", [
                "统一组织保存逻辑，接收完成后记录输入",
                "调用持久化层完成落库",
                "保存后 flag 置位，防止重复覆盖",
                "避免「已完成未保存」状态导致数据丢失",
            ])
    b._card(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(2.5),
            "🗂", "CsvSampleWriter — CSV 输出", [
                "固定路径格式：{BaseDir}/TestData/{productid}/{testid}/sensor_data.csv",
                "自动创建目录结构，稳定的文件名约定",
                "不依赖窗体对象，纯文件路径操作",
            ])

    b._add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.35),
                   "ExportCoordinator — 多格式导出", Pt(20), PRIMARY, bold=True)
    exports = [
        ("📄", "CSV 导出", "原始传感器数据\n按时间序列输出"),
        ("📊", "Excel 导出", "EPPlus 7.5.2\n格式化报表含图表"),
        ("📕", "PDF 导出", "PDFsharp-MigraDoc\n正式试验报告"),
    ]
    for i, (icon, title, desc) in enumerate(exports):
        left = Inches(0.8 + i * 4.1)
        b._card(slide, left, Inches(4.95), Inches(3.7), Inches(1.9), icon, title, [desc])

    # ================================================================
    # SLIDE 14: CHAPTER 6 — 导出亮点
    # ================================================================
    slide = b._content_slide("导出模块设计亮点", "Chapter 6 · 导出与产物输出")
    b._card(slide, Inches(0.8), Inches(1.7), Inches(3.7), Inches(2.5),
            "🔄", "Dialog 与逻辑分离", [
                "Dialog 只负责输入与触发",
                "保存逻辑在 TestRecordCoordinator",
                "导出逻辑在 ExportCoordinator",
                "文件路径在 Infrastructure/FileStorage",
            ])
    b._card(slide, Inches(4.8), Inches(1.7), Inches(3.7), Inches(2.5),
            "📦", "多格式支持", [
                "CSV：原始数据，便于二次分析",
                "Excel：格式化报表，EPPlus 生成",
                "PDF：正式报告，PDFsharp 生成",
                "一键导出，统一入口",
            ])
    b._card(slide, Inches(8.8), Inches(1.7), Inches(3.9), Inches(2.5),
            "🔒", "保存闭环", [
                "保存成功后 flag 置位",
                "防止已完成记录被覆盖",
                "未保存状态可被查询识别",
                "保证数据完整性",
            ])

    b._add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.35),
                   "导出模块架构", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(4.95), Inches(5.8), Inches(0.5), [
        "UI 层：ExportDialog 收集导出参数",
        "Features 层：ExportCoordinator 组织导出流程",
        "Infrastructure 层：文件生成与路径管理",
        "不跨层调用，不把路径写进窗体代码",
    ], Pt(13))
    b._bullet_list(slide, Inches(7.0), Inches(4.95), Inches(5.8), Inches(0.5), [
        "TestRecordCoordinator 协调保存流程",
        "CsvSampleWriter 专注 CSV 文件输出",
        "ExportCoordinator 统一多格式导出入口",
        "各司其职，单一职责原则",
    ], Pt(13))

    # ================================================================
    # SLIDE 15: CHAPTER 7 — 分工
    # ================================================================
    slide = b._content_slide("5 人分工与职责", "Chapter 7 · 团队协作与工程实践")
    roles = [
        ["角色", "负责人", "负责模块", "核心职责"],
        ["骨架/集成", "集成负责人", "sln/csproj/Shared/Config/App", "项目结构、共享契约、组合根、集成验收"],
        ["运行时", "chaoxingstar", "Core/Runtime/Shared/Events", "状态机、仿真节拍、广播契约"],
        ["UI", "Nlpexi-e", "UI/Forms/Auth/TestExecution", "主界面、弹窗、按钮矩阵、线程安全"],
        ["持久化", "userswdrs", "Persistence/History/Calibration", "SQLite、种子数据、历史查询、校准"],
        ["导出", "Y-Yunye", "FileStorage/TestRecord/Export", "试验记录、CSV、Excel/PDF 导出"],
    ]
    b._add_table(slide, Inches(0.3), Inches(1.8),
                 [Inches(1.2), Inches(1.5), Inches(3.8), Inches(5.2)],
                 roles[0], roles[1:], row_h=Inches(0.55))

    b._add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.35),
                   "5 条绝对边界规则", Pt(20), PRIMARY, bold=True)
    rules = [
        "Core 不依赖外部实现（禁止 WinForms / SQLite / 文件路径拼接）",
        "UI 不直接写业务逻辑（禁止 SQL / 状态机判断 / 导出路径拼接）",
        "Persistence 不反向依赖 UI（禁止依赖 Form / 控件对象）",
        "Export / TestRecord 不把逻辑塞进 Dialog（Dialog 只负责输入与触发）",
        "共享契约通过 Shared/Models、Shared/Events、Config/ 交互",
    ]
    for i, r in enumerate(rules):
        b._add_textbox(slide, Inches(0.8), Inches(5.95 + i * 0.28), Inches(11.5), Inches(0.25),
                       f"0{i+1}. {r}", Pt(11), TEXT_BODY)

    # ================================================================
    # SLIDE 16: CHAPTER 7 — 开发流程
    # ================================================================
    slide = b._content_slide("开发流程与工程实践", "Chapter 7 · 团队协作与工程实践")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "6 阶段推进", Pt(20), PRIMARY, bold=True)
    phases = [
        "第1阶段：打牢公共底座（骨架+运行时+持久化）",
        "第2阶段：UI 外壳 + 运行时主链路",
        "第3阶段：新建试验 + 数据落点",
        "第4阶段：完成记录 + CSV + 导出",
        "第5阶段：历史查询 + 校准",
        "第6阶段：全链路联调",
    ]
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), phases, Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "工程实践", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "Git 分支策略：每人独立 feature 分支",
        "提交前验证：build + test + run 三关",
        "公共文件变更先沟通，单独 commit",
        "5 轮联调逐步接入，不一次性全员合并",
        "验收口径：可构建 + 可测试 + 可运行",
        "冲突处理：先恢复 build，再恢复主链路",
    ], Pt(13))

    b._card(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.3),
            "✨", "协作亮点", [
                "5 人并行开发，零冲突合并 — 边界规则执行到位，每个模块有明确的文件范围和接口契约",
                "6 阶段 + 5 轮联调 — 逐步集成，降低风险，任何时候仓库都保持可构建、可测试状态",
            ])

    # ================================================================
    # SLIDE 17: CHAPTER 8 — 测试金字塔
    # ================================================================
    slide = b._content_slide("三层测试金字塔", "Chapter 8 · 测试体系")
    pyramid = [
        ("UI 自动化测试\nFlaUI", "桌面自动化验收 · TC01/TC02 用例编号",
         Inches(3.5), Inches(2.0), Inches(4.5), Inches(1.2), ACCENT),
        ("集成测试\nxUnit", "数据库初始化 · CSV 路径生成 · 导出链路",
         Inches(2.5), Inches(3.4), Inches(7.0), Inches(1.3), SECONDARY),
        ("单元测试\nxUnit", "状态机转移 · 仿真计算 · 按钮矩阵 · 边界条件 · Ready→Preparing 回退",
         Inches(1.5), Inches(4.9), Inches(9.0), Inches(1.5), PRIMARY),
    ]
    for label, desc, left, top, width, height, color in pyramid:
        b._add_rounded_rect(slide, left, top, width, height, color)
        b._add_textbox(slide, left + Inches(0.3), top + Inches(0.1), Inches(2.8), Inches(0.8),
                       label, Pt(15), WHITE, bold=True)
        b._add_textbox(slide, left + Inches(3.3), top + Inches(0.15), width - Inches(3.8),
                       Inches(0.8), desc, Pt(12), RGBColor(0xEE, 0xEE, 0xEE))

    b._add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.3),
                   "验证标准：build 通过 + test 通过 + run 启动正常 = 完成",
                   Pt(14), ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 18: CHAPTER 8 — 测试覆盖
    # ================================================================
    slide = b._content_slide("测试覆盖与验证标准", "Chapter 8 · 测试体系")
    b._card(slide, Inches(0.8), Inches(1.7), Inches(3.7), Inches(2.5),
            "🔬", "单元测试覆盖", [
                "状态转移：5 种状态 + 回退",
                "仿真初值 / 阶段推进",
                "广播快照结构验证",
                "按钮状态矩阵逻辑",
                "CSV 路径生成、记录保存参数完整性",
            ])
    b._card(slide, Inches(4.8), Inches(1.7), Inches(3.7), Inches(2.5),
            "🔗", "集成测试覆盖", [
                "数据库初始化（6 张表）",
                "种子数据写入（可重复执行）",
                "导出链路端到端",
                "历史查询多条件组合",
                "校准记录读写",
            ])
    b._card(slide, Inches(8.8), Inches(1.7), Inches(3.9), Inches(2.5),
            "🖥", "UI 自动化测试", [
                "FlaUI 桌面自动化框架",
                "Windows UI Automation",
                "TC01/TC02 等用例编号",
                "模拟真实用户操作流程",
                "验收关键业务路径",
            ])

    b._add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.35),
                   "模块验证方式", Pt(20), PRIMARY, bold=True)
    verify = [
        ["模块", "自动验证", "手动验证", "代码检查"],
        ["运行时", "dotnet test --filter Runtime", "启动验证广播", "Core 无 WinForms 引用"],
        ["UI", "dotnet build", "主窗体 + 无跨线程异常", "无 SQL / 无状态机判断"],
        ["持久化", "dotnet test --filter Persistence", "数据库文件 + 表结构", "UI 层无 SQL"],
        ["导出", "dotnet test --filter Features", "CSV 路径 + 文件输出", "Dialog 无 SQL/路径拼接"],
    ]
    b._add_table(slide, Inches(0.3), Inches(4.95),
                 [Inches(1.1), Inches(3.2), Inches(3.5), Inches(3.8)],
                 verify[0], verify[1:], row_h=Inches(0.48))

    # ================================================================
    # SLIDE 19: CHAPTER 9 — 项目成果
    # ================================================================
    slide = b._content_slide("项目成果", "Chapter 9 · 成果展示与总结")
    results = [
        "✅ 登录认证系统 — 角色选择 + 密码验证 + 错误提示",
        "✅ 主界面仪表盘 — 5通道温度显示 + 实时曲线 + 状态指示",
        "✅ 新建试验流程 — 样品信息录入 + 参数配置 + 数据落库",
        "✅ 运行时引擎 — 状态机 + PID仿真 + 800ms广播 + 线程安全",
        "✅ 试验记录保存 — 火焰现象 + 质量损失 + 自动计算 + 防重复",
        "✅ 多格式导出 — CSV 原始数据 + Excel 报表 + PDF 报告",
        "✅ 历史查询 — 按日期/样品/操作员多条件检索",
        "✅ 设备校准 — 多点校准记录 + JSON 参数 + 实时温度",
        "✅ 测试体系 — 单元测试 + 集成测试 + UI 自动化测试",
    ]
    for i, r in enumerate(results):
        col = i % 3
        row = i // 3
        b._add_textbox(slide, Inches(0.6 + col * 4.2), Inches(1.75 + row * 1.75),
                       Inches(3.9), Inches(0.35), r, Pt(13), TEXT_DARK, bold=True)

    # ================================================================
    # SLIDE 20: CHAPTER 9 — 总结
    # ================================================================
    slide = b._content_slide("技术收获与展望", "Chapter 9 · 成果展示与总结")
    b._add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.35),
                   "技术收获", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.8), Inches(0.5), [
        "分层架构实践：Core/App/Tests 三层分离",
        "状态机设计：5 状态 + 自动迁移 + 回退",
        "WinForms 线程安全：InvokeRequired 模式",
        "仿真模型：PID 控制 + 线性回归温漂计算",
        "契约驱动：RuntimeSnapshot 统一广播",
        "5 人协作：边界规则 + 接口契约 + 分阶段集成",
    ], Pt(13))

    b._add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.35),
                   "不足与改进方向", Pt(20), PRIMARY, bold=True)
    b._bullet_list(slide, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.5), [
        "DI 容器：当前手动构造，可引入 Microsoft.Extensions.DI",
        "CI/CD：缺少自动化构建流水线",
        "仿真精度：可增加更多物理模型参数",
        "报表美化：Excel/PDF 模板可进一步优化",
        "日志系统：可增加远程日志收集",
        "配置管理：可支持多环境配置切换",
    ], Pt(13))

    b._add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.3), PRIMARY)
    b._add_textbox(slide, Inches(1.2), Inches(5.45), Inches(10.5), Inches(0.4),
                   "总结", Pt(22), ACCENT, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.55),
                   "本项目实现了一个完整的材料燃烧试验仿真系统，从架构设计到编码实现，"
                   "从单元测试到 UI 自动化验收，从 5 人并行协作到全链路集成——"
                   "完整覆盖了软件开发的全生命周期，是一次从理论到实践的完整工程训练。",
                   Pt(13), WHITE, line_spacing=1.5)

    # ================================================================
    # SLIDE 21: THANK YOU
    # ================================================================
    slide = b.prs.slides.add_slide(b.prs.slide_layouts[6])
    b._set_bg(slide, PRIMARY)
    b._add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    b._add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x08, 0x1F, 0x3A))
    b._add_textbox(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(1.2),
                   "感谢聆听", Pt(52), WHITE, bold=True)
    b._add_textbox(slide, Inches(1.2), Inches(3.5), Inches(11), Inches(0.6),
                   "Thank You for Listening", Pt(24), ACCENT, font_name=EN_FONT)
    b._add_rect(slide, Inches(1.2), Inches(4.3), Inches(2.5), Inches(0.04), ACCENT)
    b._add_textbox(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(0.4),
                   "ISO 11820 材料燃烧试验仿真系统", Pt(18), RGBColor(0xAA, 0xBB, 0xCC))
    b._add_textbox(slide, Inches(1.2), Inches(5.2), Inches(11), Inches(0.4),
                   "欢迎提问与交流", Pt(16), RGBColor(0x88, 0x99, 0xAA))
    b.slide_num += 1

    # ================================================================
    # SAVE
    # ================================================================
    output_dir = r"D:\jineng\jinengshijain\jinengshijain\docs"
    output_path = os.path.join(output_dir, "ISO11820-课程项目答辩-v2.pptx")
    b.prs.save(output_path)
    print(f"PPT saved to: {output_path}")
    print(f"Total slides: {b.slide_num}")


if __name__ == "__main__":
    build()